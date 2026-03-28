"""
GFD LLM Slide Generator
=======================
Two-stage pipeline:

Stage 1 — LLM slide spec:
  Receives the structured JSON from gfd_llm_parser and generates a complete
  slide specification — RAG colors, condensed cell text, groupings, overview
  bullets, risk assessments. The LLM makes ALL content and layout decisions.

Stage 2 — PPTX renderer:
  A thin "paint by numbers" renderer that converts the slide spec JSON into a
  python-pptx Presentation object. No business logic lives here — the renderer
  just maps JSON fields to shapes, colors, and text runs.

If Stage 1 fails, a secondary LLM fallback processes one product group at a time.
If an individual per-PG call also fails, a skeleton slide (blank text, GREY cells)
is used for that PG only — ensuring the renderer never crashes.
"""

from __future__ import annotations

import io
import json
import time
from datetime import date, datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage

from agent import log_tokens, log_trace


# ─── LLM factory (local, higher max_tokens for large JSON output) ─────

def _create_llm(config: dict, max_tokens: int = 8192) -> AzureChatOpenAI:
    return AzureChatOpenAI(
        azure_deployment=config["azure_deployment"],
        azure_endpoint=config["azure_endpoint"],
        api_key=config["api_key"],
        api_version=config.get("api_version", "2024-12-01-preview"),
        temperature=0.15,
        max_tokens=max_tokens,
    )


# ─── Color palette ────────────────────────────────────────────────────

_C: dict[str, RGBColor] = {
    "GREEN":      RGBColor(0x00, 0xB0, 0x50),
    "AMBER":      RGBColor(0xFF, 0xC0, 0x00),
    "RED":        RGBColor(0xC0, 0x00, 0x00),
    "GREY":       RGBColor(0xA0, 0xA0, 0xA0),
    "WHITE":      RGBColor(0xFF, 0xFF, 0xFF),
    "BLACK":      RGBColor(0x1A, 0x1A, 0x1A),
    "NAVY":       RGBColor(0x1E, 0x27, 0x61),
    "BLUE":       RGBColor(0x2C, 0x52, 0x82),
    "LIGHT_BLUE": RGBColor(0xBF, 0xD7, 0xED),
    "LIGHT_GREY": RGBColor(0xF4, 0xF4, 0xF6),
    "MEDIUM_ORANGE": RGBColor(0xFF, 0x99, 0x00),
}

_RAG_FG: dict[str, RGBColor] = {
    "GREEN": _C["WHITE"],
    "AMBER": _C["BLACK"],
    "RED":   _C["WHITE"],
    "GREY":  _C["WHITE"],
}

_RISK_BG: dict[str, RGBColor] = {
    "CRITICAL": _C["RED"],
    "HIGH":     _C["AMBER"],
    "MEDIUM":   _C["MEDIUM_ORANGE"],
    "LOW":      _C["GREEN"],
}
_RISK_FG: dict[str, RGBColor] = {
    "CRITICAL": _C["WHITE"],
    "HIGH":     _C["BLACK"],
    "MEDIUM":   _C["BLACK"],
    "LOW":      _C["WHITE"],
}

# Slide dimensions — widescreen 16:9
_W  = Inches(13.333)
_H  = Inches(7.5)
_M  = Inches(0.28)     # side margin
_HH = Inches(0.62)     # header bar height

FONT = "Calibri"


# ─── Low-level python-pptx helpers ───────────────────────────────────

def _solid(shape, rgb: RGBColor) -> None:
    """Solid fill + no border on a shape."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _rect(slide, x, y, w, h, rgb: RGBColor):
    """Add a filled rectangle (MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1)."""
    s = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    _solid(s, rgb)
    return s


def _tb(slide, x, y, w, h, text: str, *,
        size: int = 10, bold: bool = False,
        color: RGBColor = None, align=PP_ALIGN.LEFT,
        bg: RGBColor = None) -> Any:
    """Add a textbox. bg=None → transparent fill."""
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    if bg:
        tb.fill.solid()
        tb.fill.fore_color.rgb = bg
    else:
        tb.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def _cell_fill(cell, rgb: RGBColor) -> None:
    """Set solid background fill on a table cell via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing fill elements
    for child in list(tcPr):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            tcPr.remove(child)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(sf, qn("a:srgbClr"))
    clr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _cell_write(cell, text: str, *,
                size: int = 7, bold: bool = False,
                color: RGBColor = None, align=PP_ALIGN.LEFT) -> None:
    """Write text into a table cell, clearing previous content."""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    # Remove extra runs from first paragraph
    for r in list(p.runs):
        p._p.remove(r._r)
    # Remove extra paragraphs
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


# ─── Stage 1: LLM slide spec ─────────────────────────────────────────

_SLIDES_SYSTEM = """\
You are a Chief Supply Chain Officer designing a Global Fulfilment Dashboard PowerPoint
presentation for a board-level audience. You are given JSON extracted from the
Dashboard_Update worksheet. Generate a COMPLETE slide specification in JSON.

═══ SLIDE DESIGN RULES ═══

SLIDE 1 — Executive Overview (type: "overview")
  • overall_risk: worst-case across ALL rows (CRITICAL > HIGH > MEDIUM > LOW).
  • 4–6 bullet points (max 110 chars each). Lead with the most critical items.
    Every bullet must cite specific facts: CW numbers, plant names, counts.
  • stats: count risk items by risk level.

SLIDES 2+ — One slide per product group (type: "product_group")
  • If a product group has >9 rows, split it into multiple slides with a
    "(cont.)" suffix on the title.
  • title: "Product Group Desc (CODE)" — use actual codes and descriptions.
  • headline: 1-sentence situation summary (max 80 chars). Be specific.
  • cw_columns: list of exactly 12 consecutive CW integers starting from
    the current CW number (e.g. if current is CW13 → [13,14,15,16,17,18,19,20,21,22,23,24]).
    Handle year-end wrap correctly (e.g. CW51,52,1,2,...).
  • quarter_label: next quarter label (e.g. "Q2/2026").

RAG COLOR COMPUTATION (per row, per CW column):
  Given coverage_without_mitigation_cw (W) and coverage_with_mitigation_cw (M):
    • "GREEN" if CW ≤ W  (supply secured without mitigation)
    • "AMBER" if W < CW ≤ M  (supply depends on mitigation actions)
    • "RED"   if CW > M  (no supply plan in place)
    • If M is null but W is set: "GREEN" if CW ≤ W, else "RED"
    • If both are null: "GREY"
  quarter_color = worst RAG across all 13 weeks of that quarter
  (RED > AMBER > GREEN > GREY).

KB COVERAGE FIELD (per row):
  kb_coverage = the furthest CW through which supply to the customer is secured,
  formatted as "CW{week}/{2-digit year}" (e.g. "CW19/26").
  Use coverage_with_mitigation_cw if available, else coverage_without_mitigation_cw.
  If both are null use "N/A".

TEXT CONDENSING for slide display (applied to each row field):
  product_group ≤ 30 chars (PG code + short name, e.g. "SEN - Sensors"),
  plant ≤ 15 chars, customer ≤ 40 chars, supplier ≤ 30 chars,
  action ≤ 90 chars, fm_status ≤ 15 chars.

═══ OUTPUT FORMAT ═══

Respond with ONLY valid JSON — no markdown fences, no explanation:

{{
  "presentation_title": "Global Fulfilment Dashboard",
  "current_cw": "CW13/2026",
  "generated_date": "{today}",
  "slides": [
    {{
      "type": "overview",
      "title": "Situation Overview — CW13/2026",
      "overall_risk": "HIGH",
      "bullets": ["3 plants at RED from CW16 — BMW Group coverage expires CW15", "..."],
      "stats": {{
        "total_items": 12,
        "critical_count": 2,
        "high_count": 4,
        "medium_count": 4,
        "low_count": 2,
        "product_groups_count": 3
      }}
    }},
    {{
      "type": "product_group",
      "title": "Sensors / Radar (SEN)",
      "headline": "NXP wafer shortage — BHV at RED from CW16, mitigation active",
      "cw_columns": [13,14,15,16,17,18,19,20,21,22,23,24],
      "quarter_label": "Q2/2026",
      "rows": [
        {{
          "product_group": "SEN - Sensors",
          "plant": "BHV",
          "customer": "BMW Group, Mercedes-Benz",
          "kb_coverage": "CW19/26",
          "supplier": "NXP Semiconductors",
          "coverage_without": 15,
          "coverage_with": 19,
          "action": "Dual source by CW14; air freight bridge until CW18",
          "fm_status": "In progress",
          "cw_colors": {{"13":"GREEN","14":"GREEN","15":"GREEN","16":"AMBER","17":"AMBER","18":"AMBER","19":"AMBER","20":"RED","21":"RED","22":"RED","23":"RED","24":"RED"}},
          "quarter_color": "RED"
        }}
      ]
    }}
  ]
}}
{glossary_block}"""

_SLIDES_USER = """\
Today: {today}
Current calendar week: {current_cw}

Generate the complete dashboard slide specification from this extracted data:

{extracted_json}"""


async def llm_generate_slide_spec(
    extracted: dict,
    llm_config: dict,
    session_id: str,
    glossary_context: str = "",
) -> dict:
    """
    Stage 3: LLM generates a complete slide specification from extracted data.

    Primary path  — one call, full spec for all product groups.
    Secondary path — if primary fails (truncation, parse error, API error):
                     one LLM call per product group + one overview call.
    Last resort   — if a per-PG secondary call also fails: skeleton slide
                    for that PG only (blank text, GREY RAG cells).
    No deterministic RAG computation anywhere in this stack.
    """
    llm = _create_llm(llm_config, max_tokens=8192)
    t0 = time.time()
    today = date.today().isoformat()
    current_cw = extracted.get("current_cw", "CW??/????")

    glossary_block = (
        f"\n\nCOMPANY GLOSSARY:\n{glossary_context}" if glossary_context else ""
    )

    # Exclude internal metadata keys from what the LLM sees
    clean = {k: v for k, v in extracted.items() if not k.startswith("_")}
    extracted_json = json.dumps(clean, indent=2, ensure_ascii=False)

    messages = [
        SystemMessage(content=_SLIDES_SYSTEM.format(
            today=today,
            glossary_block=glossary_block,
        )),
        HumanMessage(content=_SLIDES_USER.format(
            today=today,
            current_cw=current_cw,
            extracted_json=extracted_json,
        )),
    ]

    try:
        response = await llm.ainvoke(messages)
        raw = response.content.strip()

        if raw.startswith("```"):
            raw = "\n".join(raw.split("\n")[1:])
        if raw.endswith("```"):
            raw = "\n".join(raw.split("\n")[:-1])

        spec: dict = json.loads(raw.strip())

        usage = response.response_metadata.get("token_usage", {})
        log_tokens(session_id, "gfd_llm_slide_spec", usage, llm_config.get("azure_deployment", ""))

        n_slides = len(spec.get("slides", []))
        duration = (time.time() - t0) * 1000
        log_trace(
            session_id, "gfd_llm_slide_spec",
            f"Input: {len(extracted.get('product_groups', []))} product groups",
            f"Generated {n_slides} slide specs",
            duration,
        )
        return spec

    except Exception as exc:
        # Primary call failed — log and try per-PG secondary LLM fallback
        duration = (time.time() - t0) * 1000
        log_trace(session_id, "gfd_llm_slide_spec",
                  "Primary slide spec call",
                  f"FAILED ({type(exc).__name__}: {str(exc)[:100]}) — trying per-PG fallback",
                  duration, {"error": True, "fallback": True})
        return await _llm_fallback_spec(extracted, llm_config, session_id, glossary_context)


# ─── Secondary LLM fallback ───────────────────────────────────────────
# Called when the primary (full-spec) LLM call fails.
# Processes one product group at a time — smaller calls are less likely
# to be truncated and failures are isolated per PG.

_FALLBACK_PG_SYSTEM = """\
You are generating a single slide specification for a Global Fulfilment Dashboard.

You will receive raw extracted data for ONE product group. Return ONLY valid JSON
for that one slide — no markdown fences, no explanation.

CURRENT CALENDAR WEEK: {current_cw}

CW GRID: produce exactly 12 consecutive CW integers starting from the current CW,
wrapping at 52 (e.g. CW51 → CW52 → CW1 → CW2 ...).

RAG COLORS per cell:
  "GREEN" if CW ≤ coverage_without_mitigation_cw
  "AMBER" if coverage_without < CW ≤ coverage_with_mitigation_cw
  "RED"   if CW > coverage_with_mitigation_cw (or both null → "GREY")

KB COVERAGE: use coverage_with_mitigation_cw if set, else coverage_without_mitigation_cw.
Format as "CW{{week}}/{{2-digit year}}" (e.g. "CW19/26"). Use "N/A" if both are null.

PRODUCT GROUP field: combine code + short description, max 30 chars (e.g. "SEN - Sensors").
TEXT LIMITS: plant ≤ 15, customer ≤ 40, supplier ≤ 30, action ≤ 90, fm_status ≤ 15 chars.

OUTPUT SCHEMA — exactly this structure:
{{
  "type": "product_group",
  "title": "<Product Family Desc> (<CODE>)",
  "headline": "<1-sentence situation summary, max 80 chars>",
  "cw_columns": [<12 consecutive CW integers>],
  "quarter_label": "<e.g. Q2/2026>",
  "rows": [
    {{
      "product_group": "<code - short name>",
      "plant": "<plant code>",
      "customer": "<affected customers>",
      "kb_coverage": "<CW##/YY or N/A>",
      "supplier": "<supplier name>",
      "action": "<condensed action/comment>",
      "fm_status": "<Yes|No|In progress|N/A>",
      "cw_colors": {{"<cw>": "<GREEN|AMBER|RED|GREY>", ...}},
      "quarter_color": "<GREEN|AMBER|RED|GREY>"
    }}
  ]
}}
{glossary_block}"""

_FALLBACK_PG_USER = """\
Product group data (current CW: {current_cw}):

{pg_json}"""

_FALLBACK_OVERVIEW_SYSTEM = """\
You are generating the executive overview slide for a Global Fulfilment Dashboard.

You are given a list of product-group slide specs that have already been generated.
Write the overview slide JSON — no markdown fences, no explanation.

OUTPUT SCHEMA:
{{
  "type": "overview",
  "title": "Situation Overview — {current_cw}",
  "overall_risk": "<CRITICAL|HIGH|MEDIUM|LOW>",
  "bullets": ["<up to 6 bullets, max 110 chars each, cite CW numbers and plant names>"],
  "stats": {{
    "total_items": <int>,
    "critical_count": <int>,
    "high_count": <int>,
    "medium_count": <int>,
    "low_count": <int>,
    "product_groups_count": <int>
  }}
}}

Derive overall_risk as worst across all rows (RED→CRITICAL, any AMBER→HIGH, all GREEN→LOW).
Derive stats by counting rows whose cw_colors contain at least one RED (CRITICAL),
at least one AMBER but no RED (HIGH), all GREEN (LOW), else MEDIUM.
{glossary_block}"""

_FALLBACK_OVERVIEW_USER = """\
Current CW: {current_cw}
Total product groups: {n_pgs}

Product group slide specs:
{slides_json}"""


async def _llm_fallback_spec(
    extracted: dict,
    llm_config: dict,
    session_id: str,
    glossary_context: str = "",
) -> dict:
    """
    Secondary LLM fallback: one LLM call per product group, then one overview call.
    Called when the primary full-spec call fails (truncation, API error, etc.).
    Failures on individual PGs fall back to a skeleton slide for that PG only.
    """
    llm = _create_llm(llm_config, max_tokens=4096)
    current_cw = extracted.get("current_cw", "CW??/????")
    today = date.today().isoformat()

    glossary_block = (
        f"\n\nCOMPANY GLOSSARY:\n{glossary_context}" if glossary_context else ""
    )

    pg_slides: list[dict] = []

    for pg_idx, pg in enumerate(extracted.get("product_groups", [])):
        t0 = time.time()
        pg_name = pg.get("product_family_desc", "Unknown")
        pg_code = pg.get("product_family_code", "")
        pg_display = f"{pg_name} ({pg_code})" if pg_code else pg_name
        pg_json = json.dumps(pg, indent=2, ensure_ascii=False)

        messages = [
            SystemMessage(content=_FALLBACK_PG_SYSTEM.format(
                current_cw=current_cw,
                glossary_block=glossary_block,
            )),
            HumanMessage(content=_FALLBACK_PG_USER.format(
                current_cw=current_cw,
                pg_json=pg_json,
            )),
        ]

        try:
            response = await llm.ainvoke(messages)
            raw = response.content.strip()
            if raw.startswith("```"):
                raw = "\n".join(raw.split("\n")[1:])
            if raw.endswith("```"):
                raw = "\n".join(raw.split("\n")[:-1])

            pg_slide: dict = json.loads(raw.strip())

            usage = response.response_metadata.get("token_usage", {})
            log_tokens(session_id, f"gfd_fallback_pg_{pg_idx}", usage,
                       llm_config.get("azure_deployment", ""))
            log_trace(session_id, "gfd_fallback_pg",
                      f"PG: {pg_display}",
                      pg_slide.get("headline", "")[:120],
                      (time.time() - t0) * 1000, {"pg_index": pg_idx})

            pg_slides.append(pg_slide)

        except Exception as exc:
            # Individual PG failure → skeleton slide for this PG only
            log_trace(session_id, "gfd_fallback_pg",
                      f"PG: {pg_display}",
                      f"SKELETON (error: {str(exc)[:80]})",
                      (time.time() - t0) * 1000, {"error": True, "pg_index": pg_idx})
            pg_slides.append(_skeleton_pg_slide(pg, current_cw))

    # ── Overview slide ────────────────────────────────────────────────
    t0 = time.time()
    try:
        overview_messages = [
            SystemMessage(content=_FALLBACK_OVERVIEW_SYSTEM.format(
                current_cw=current_cw,
                glossary_block=glossary_block,
            )),
            HumanMessage(content=_FALLBACK_OVERVIEW_USER.format(
                current_cw=current_cw,
                n_pgs=len(pg_slides),
                slides_json=json.dumps(pg_slides, indent=2, ensure_ascii=False),
            )),
        ]
        ov_response = await llm.ainvoke(overview_messages)
        raw_ov = ov_response.content.strip()
        if raw_ov.startswith("```"):
            raw_ov = "\n".join(raw_ov.split("\n")[1:])
        if raw_ov.endswith("```"):
            raw_ov = "\n".join(raw_ov.split("\n")[:-1])

        overview_slide = json.loads(raw_ov.strip())

        usage = ov_response.response_metadata.get("token_usage", {})
        log_tokens(session_id, "gfd_fallback_overview", usage,
                   llm_config.get("azure_deployment", ""))
        log_trace(session_id, "gfd_fallback_overview",
                  f"Overview for {len(pg_slides)} PG slides",
                  overview_slide.get("overall_risk", "")[:50],
                  (time.time() - t0) * 1000)

    except Exception as exc:
        # Overview LLM call failed → minimal skeleton overview
        log_trace(session_id, "gfd_fallback_overview",
                  "Overview generation",
                  f"SKELETON OVERVIEW (error: {str(exc)[:80]})",
                  (time.time() - t0) * 1000, {"error": True})
        overview_slide = _skeleton_overview(extracted, current_cw)

    return {
        "presentation_title": "Global Fulfilment Dashboard",
        "current_cw":         current_cw,
        "generated_date":     today,
        "_fallback":          True,
        "slides":             [overview_slide] + pg_slides,
    }


# ─── Skeleton spec — last resort only ────────────────────────────────
# Used only when an individual LLM call within the fallback itself fails.
# Produces structurally valid slides with blank/grey cells — no content
# computation whatsoever. The renderer will never crash on this output.

def _parse_cw_start(current_cw_str: str) -> tuple[int, int]:
    """Extract (cw_week, cw_year) from a 'CW13/2026' string, defaulting to today."""
    import re as _re
    from datetime import datetime as _dt
    m = _re.match(r"CW(\d+)/(\d+)", current_cw_str or "")
    if m:
        return int(m.group(1)), int(m.group(2))
    iso = _dt.now().isocalendar()
    return iso.week, iso.year


def _skeleton_cw_columns(cw_start: int) -> list[int]:
    return [((cw_start - 1 + i) % 52) + 1 for i in range(12)]


def _skeleton_pg_slide(pg: dict, current_cw_str: str) -> dict:
    """Minimal skeleton for one product group — blank text, GREY RAG cells."""
    code = pg.get("product_family_code", "")
    desc = pg.get("product_family_desc", "Unknown")
    title = f"{desc} ({code})" if code else desc

    cw_start, cw_year = _parse_cw_start(current_cw_str)
    cw_columns = _skeleton_cw_columns(cw_start)
    q_num = min(((cw_columns[-1] - 1) // 13) + 2, 4)
    q_label = f"Q{q_num}/{cw_year}"

    rows = []
    for row in pg.get("rows", []):
        pg_label = f"{code} - {desc[:18]}" if code else desc[:28]
        cw_colors = {str(cw): "GREY" for cw in cw_columns}
        rows.append({
            "product_group": pg_label[:30],
            "plant":         (row.get("plant_location") or "")[:15],
            "customer":      (row.get("customer_affected") or "")[:40],
            "kb_coverage":   "N/A",
            "supplier":      (row.get("supplier_text") or "")[:30],
            "action":        (row.get("action_comment") or "")[:90],
            "fm_status":     str(row.get("customer_informed") or "N/A")[:15],
            "cw_colors":     cw_colors,
            "quarter_color": "GREY",
        })

    return {
        "type":          "product_group",
        "title":         title,
        "headline":      "Data available — RAG computation pending (reprocess to generate)",
        "cw_columns":    cw_columns,
        "quarter_label": q_label,
        "rows":          rows,
    }


def _skeleton_overview(extracted: dict, current_cw_str: str) -> dict:
    """Minimal skeleton overview when the overview LLM call itself fails."""
    n_pgs = len(extracted.get("product_groups", []))
    total = sum(len(pg.get("rows", [])) for pg in extracted.get("product_groups", []))
    return {
        "type":         "overview",
        "title":        f"Situation Overview — {current_cw_str}",
        "overall_risk": "HIGH",
        "bullets": [
            f"{n_pgs} product group(s) with active fulfilment risks",
            f"{total} risk items tracked — see per-group slides for detail",
            "Reprocess to generate AI-driven risk narrative and RAG analysis",
        ],
        "stats": {
            "total_items":          total,
            "critical_count":       0,
            "high_count":           0,
            "medium_count":         total,
            "low_count":            0,
            "product_groups_count": n_pgs,
        },
    }


# ─── Stage 2: PPTX renderer ──────────────────────────────────────────

def _render_overview(prs: Presentation, spec: dict) -> None:
    """Render the executive overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    title        = spec.get("title", "Global Fulfilment Dashboard")
    overall_risk = spec.get("overall_risk", "HIGH").upper()
    bullets      = spec.get("bullets", [])
    stats        = spec.get("stats", {})

    # ── Header bar ────────────────────────────────────────────────────
    _rect(slide, 0, 0, _W, _HH, _C["NAVY"])
    _tb(slide, _M, Inches(0.1), Inches(10.5), Inches(0.45),
        title, size=20, bold=True, color=_C["WHITE"])

    # Risk badge (top-right corner of header)
    badge_w = Inches(1.75)
    badge_x = _W - badge_w - _M
    _rect(slide, badge_x, Inches(0.09), badge_w, Inches(0.44),
          _RISK_BG.get(overall_risk, _C["AMBER"]))
    _tb(slide, badge_x, Inches(0.09), badge_w, Inches(0.44),
        f"● {overall_risk} RISK",
        size=11, bold=True,
        color=_RISK_FG.get(overall_risk, _C["BLACK"]),
        align=PP_ALIGN.CENTER)

    # ── Stats bar ─────────────────────────────────────────────────────
    bar_y  = _HH + Inches(0.1)
    bar_h  = Inches(0.68)
    stat_w = Inches(2.38)
    sx     = _M

    for key, label, bg, fg in [
        ("critical_count", "CRITICAL", _C["RED"],          _C["WHITE"]),
        ("high_count",     "HIGH",     _C["AMBER"],        _C["BLACK"]),
        ("medium_count",   "MEDIUM",   _C["MEDIUM_ORANGE"],_C["BLACK"]),
        ("low_count",      "LOW",      _C["GREEN"],        _C["WHITE"]),
    ]:
        count = stats.get(key, 0)
        _rect(slide, sx, bar_y, stat_w, bar_h, bg)
        _tb(slide, sx, bar_y, stat_w, bar_h,
            f"{count}  {label}",
            size=13, bold=True, color=fg, align=PP_ALIGN.CENTER)
        sx += stat_w + Inches(0.12)

    # Item / PG summary (right-aligned in stats bar)
    total = stats.get("total_items", 0)
    pgs   = stats.get("product_groups_count", 0)
    _tb(slide, Inches(10.4), bar_y + Inches(0.19), Inches(2.6), Inches(0.32),
        f"{total} items  ·  {pgs} product groups",
        size=8, color=_C["BLACK"], align=PP_ALIGN.RIGHT)

    # ── Section label ─────────────────────────────────────────────────
    sub_y = bar_y + bar_h + Inches(0.14)
    _tb(slide, _M, sub_y, Inches(12.5), Inches(0.26),
        "KEY RISK HIGHLIGHTS",
        size=8, bold=True, color=_C["BLUE"])
    sub_y += Inches(0.29)

    # ── Bullets ───────────────────────────────────────────────────────
    remaining = _H - sub_y - _M
    per_bullet = min(remaining / max(len(bullets), 1), Inches(0.65))
    for bullet in bullets:
        _tb(slide, Inches(0.55), sub_y, Inches(12.1), per_bullet,
            f"▸  {bullet}", size=11, color=_C["BLACK"])
        sub_y += per_bullet

    # ── Footer ────────────────────────────────────────────────────────
    gen_date = date.today().isoformat()
    _tb(slide, Inches(10.0), _H - Inches(0.28), Inches(3.0), Inches(0.24),
        f"Generated {gen_date}", size=7,
        color=_C["GREY"], align=PP_ALIGN.RIGHT)


def _render_product_group(prs: Presentation, spec: dict) -> None:
    """Render a single product-group RAG heatmap slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title         = spec.get("title", "Product Group")
    headline      = spec.get("headline", "")
    cw_columns: list[int] = spec.get("cw_columns", [])
    quarter_label = spec.get("quarter_label", "Q?")
    rows: list[dict] = spec.get("rows", [])

    # ── Header ────────────────────────────────────────────────────────
    _rect(slide, 0, 0, _W, _HH, _C["NAVY"])
    _tb(slide, _M, Inches(0.06), Inches(11.5), Inches(0.30),
        title, size=14, bold=True, color=_C["WHITE"])
    _tb(slide, _M, Inches(0.36), Inches(11.5), Inches(0.22),
        headline, size=8, color=_C["LIGHT_BLUE"])

    if not rows:
        _tb(slide, _M, Inches(1.2), Inches(12.0), Inches(0.4),
            "No data rows for this product group.", size=11, color=_C["BLACK"])
        return

    # ── Table geometry ────────────────────────────────────────────────
    # Fixed columns: PG | Plant | Customer | KB Coverage | Supplier | Action | FM
    # Risk column removed — RAG colours convey risk implicitly.
    FIXED_LABELS = [
        "Product Group\n(PG)", "Plant",      "Customer",
        "KB Coverage\n(CW/YY)",              "Supplier",
        "Action / Comment",                  "FM",
    ]
    FIXED_WIDTHS = [
        Inches(1.40), Inches(0.62), Inches(1.12),
        Inches(0.75),              Inches(0.90),
        Inches(2.20),                         Inches(0.44),
    ]

    CW_W = Inches(0.365)
    Q_W  = Inches(0.53)

    n_cw = len(cw_columns)
    total_fixed = sum(FIXED_WIDTHS)
    total_cw    = n_cw * CW_W + Q_W
    available   = _W - 2 * _M

    # Scale CW columns proportionally if they don't fit
    if total_fixed + total_cw > available:
        scale = (available - total_fixed) / total_cw
        CW_W = int(CW_W * scale)
        Q_W  = int(Q_W  * scale)

    all_widths = FIXED_WIDTHS + [CW_W] * n_cw + [Q_W]
    n_cols     = len(all_widths)
    n_rows_tbl = len(rows) + 1  # +1 for header

    tbl_y  = _HH + Inches(0.1)
    tbl_h  = _H - tbl_y - Inches(0.28)   # space for legend
    row_h  = max(int(tbl_h / n_rows_tbl), int(Inches(0.31)))

    tbl_shape = slide.shapes.add_table(
        n_rows_tbl, n_cols,
        int(_M), int(tbl_y),
        sum(int(w) for w in all_widths),
        row_h * n_rows_tbl,
    )
    tbl = tbl_shape.table

    # Apply column widths and row heights
    for i, w in enumerate(all_widths):
        tbl.columns[i].width = int(w)
    for r in range(n_rows_tbl):
        tbl.rows[r].height = row_h

    # ── Header row ────────────────────────────────────────────────────
    hdr_labels = FIXED_LABELS + [f"CW{c}" for c in cw_columns] + [quarter_label]
    for ci, label in enumerate(hdr_labels):
        cell = tbl.cell(0, ci)
        _cell_fill(cell, _C["NAVY"])
        is_cw_col = ci >= len(FIXED_LABELS)
        _cell_write(cell, label, size=7, bold=True, color=_C["WHITE"],
                    align=PP_ALIGN.CENTER if is_cw_col else PP_ALIGN.LEFT)

    # ── Data rows ─────────────────────────────────────────────────────
    for ri, row in enumerate(rows, start=1):
        alt = _C["LIGHT_GREY"] if ri % 2 == 0 else _C["WHITE"]

        # Fixed text columns  (order must match FIXED_LABELS exactly)
        fixed_vals = [
            row.get("product_group", ""),
            row.get("plant", ""),
            row.get("customer", ""),
            row.get("kb_coverage", ""),
            row.get("supplier", ""),
            row.get("action", ""),
            row.get("fm_status", ""),
        ]
        for ci, val in enumerate(fixed_vals):
            cell = tbl.cell(ri, ci)
            _cell_fill(cell, alt)
            # Centre: Plant (1), KB Coverage (3), FM (6)
            center = ci in (1, 3, 6)
            _cell_write(cell, str(val) if val else "", size=7,
                        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
                        color=_C["BLACK"])

        # CW colored cells
        cw_colors: dict = row.get("cw_colors", {})
        for cw_i, cw_num in enumerate(cw_columns):
            ci = len(FIXED_LABELS) + cw_i
            rag = str(cw_colors.get(str(cw_num), "GREY")).upper()
            cell = tbl.cell(ri, ci)
            _cell_fill(cell, _C.get(rag, _C["GREY"]))
            _cell_write(cell, "", size=5)   # empty — color carries the meaning

        # Quarter column
        q_ci  = len(FIXED_LABELS) + n_cw
        q_rag = str(row.get("quarter_color", "GREY")).upper()
        q_cell = tbl.cell(ri, q_ci)
        _cell_fill(q_cell, _C.get(q_rag, _C["GREY"]))
        _cell_write(q_cell, "", size=5)

    # ── Color legend ──────────────────────────────────────────────────
    legend_y = tbl_y + row_h * n_rows_tbl + Inches(0.07)
    if legend_y + Inches(0.2) <= _H:
        lx = _M
        for label, color in [
            ("■ Covered (w/o mitigation)", _C["GREEN"]),
            ("■ Covered w/ mitigation",    _C["AMBER"]),
            ("■ Uncovered",                _C["RED"]),
            ("■ Data unavailable",         _C["GREY"]),
        ]:
            _tb(slide, lx, legend_y, Inches(2.85), Inches(0.2),
                label, size=6, color=color)
            lx += Inches(2.9)


def render_pptx_from_spec(slide_spec: dict, output_path: str | None = None) -> io.BytesIO:
    """
    Stage 2 (renderer): Convert an LLM-generated slide spec dict
    into a python-pptx Presentation.

    Parameters
    ----------
    slide_spec  : dict produced by llm_generate_slide_spec (or the fallback)
    output_path : optional path to save the file to disk

    Returns
    -------
    io.BytesIO buffer (seeked to 0) containing the .pptx bytes
    """
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H

    for slide_def in slide_spec.get("slides", []):
        stype = slide_def.get("type", "")
        if stype == "overview":
            _render_overview(prs, slide_def)
        elif stype == "product_group":
            _render_product_group(prs, slide_def)
        # Silently skip unknown types

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    if output_path:
        with open(output_path, "wb") as fh:
            fh.write(buf.read())
        buf.seek(0)

    return buf


# ─── Combined entry point ─────────────────────────────────────────────

async def generate_gfd_dashboard(
    extracted_data: dict,
    llm_config: dict,
    session_id: str,
    output_path: str | None = None,
    glossary_context: str = "",
) -> tuple[io.BytesIO, dict]:
    """
    Full GFD dashboard generation pipeline.

    1. LLM generates a complete slide specification from extracted_data.
    2. Thin PPTX renderer converts the spec to a python-pptx file.

    Parameters
    ----------
    extracted_data : output of parse_gfd_with_llm()
    llm_config     : Azure OpenAI config dict
    session_id     : for token / trace logging
    output_path    : optional path to save the .pptx file
    glossary_context : rendered glossary string for LLM context

    Returns
    -------
    (buf: BytesIO, slide_spec: dict)
    """
    spec = await llm_generate_slide_spec(
        extracted_data, llm_config, session_id, glossary_context
    )
    buf = render_pptx_from_spec(spec, output_path=output_path)
    return buf, spec
