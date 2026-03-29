"""
SID LLM Slide Generator
========================
Two-stage pipeline:

Stage 3 — LLM slide spec:
  Receives the structured JSON from sid_llm_parser and generates a complete
  two-slide specification:
    Slide 1: Executive overview — summary, coverage pie chart data,
             affected suppliers table, action tracking table
    Slide 2: Supplier-level fulfillment & recovery detail table

Stage 4 — PPTX renderer:
  A thin "paint by numbers" renderer matching the corporate template:
  - 10 × 5.625 inch slides (16:9)
  - Slide 1: title, pie chart legend, affected suppliers table (5 cols),
    actions table (4 cols), contextual notes, footer
  - Slide 2: 9-column supplier detail table, footer

If Stage 3 fails, a deterministic fallback builds the spec from raw data.

When a corporate template PPTX is available, it is loaded as the base.
"""

from __future__ import annotations

import io
import json
import math
import time
from datetime import date, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn
from lxml import etree

from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage

from agent import log_tokens, log_trace
from sid_llm_parser import _parse_llm_json


# ─── Template path ───────────────────────────────────────────────────

TEMPLATE_PATH = Path(__file__).parent / "sid_ppt_template.pptx"

# ─── LLM factory ────────────────────────────────────────────────────

def _create_llm(config: dict, max_tokens: int = 16384) -> AzureChatOpenAI:
    return AzureChatOpenAI(
        azure_deployment=config["azure_deployment"],
        azure_endpoint=config["azure_endpoint"],
        api_key=config["api_key"],
        api_version=config.get("api_version", "2024-12-01-preview"),
        temperature=0.15,
        max_tokens=max_tokens,
    )


# ─── Colour palette (matching the corporate template) ───────────────

_C: dict[str, RGBColor] = {
    "HDR_BLUE":     RGBColor(0x3F, 0x7A, 0xB6),   # accent1 — table headers
    "DARK_BLUE":    RGBColor(0x00, 0x45, 0x7E),   # dk2/tx2 — slide2 table header bg
    "GREEN":        RGBColor(0x38, 0xDF, 0x12),   # coverage > 15 days
    "YELLOW":       RGBColor(0xFC, 0xEF, 0x39),   # coverage 5-15 days
    "RED":          RGBColor(0xC1, 0x00, 0x1F),   # coverage < 4 days (accent5)
    "BLACK":        RGBColor(0x00, 0x00, 0x00),   # no coverage
    "WHITE":        RGBColor(0xFF, 0xFF, 0xFF),
    "GREY":         RGBColor(0xA0, 0xA0, 0xA0),
    "LIGHT_GREY":   RGBColor(0xF0, 0xF0, 0xF0),
    "BORDER":       RGBColor(0xCC, 0xCC, 0xCC),
    "SEVERITY_R":   RGBColor(0xFF, 0x00, 0x00),
    "SEVERITY_Y":   RGBColor(0xFF, 0xC0, 0x00),
    "SEVERITY_G":   RGBColor(0x00, 0xB0, 0x50),
}

# Coverage chart colors
_COVERAGE_COLORS = {
    "no_coverage":  _C["BLACK"],
    "lt_4_days":    _C["RED"],
    "5_to_15_days": _C["YELLOW"],
    "gt_15_days":   _C["GREEN"],
}

# Slide dimensions — standard 16:9
_W  = Inches(10.0)
_H  = Inches(5.625)

FONT = "Arial"


# ─── Low-level python-pptx helpers ──────────────────────────────────

def _cell_fill(cell, rgb: RGBColor) -> None:
    """Set solid background fill on a table cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for child in list(tcPr):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            tcPr.remove(child)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(sf, qn("a:srgbClr"))
    clr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _cell_margins(cell, left=36000, right=36000, top=18000, bottom=18000,
                  anchor: str = "t") -> None:
    """Set cell margins and vertical anchor."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set("marL", str(left))
    tcPr.set("marR", str(right))
    tcPr.set("marT", str(top))
    tcPr.set("marB", str(bottom))
    tcPr.set("anchor", anchor)


def _para_spacing_zero(p) -> None:
    """Set paragraph spacing to 0."""
    pPr = p._p.get_or_add_pPr()
    for old in pPr.findall(qn("a:spcBef")):
        pPr.remove(old)
    sb = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(sb, qn("a:spcPts")).set("val", "0")
    for old in pPr.findall(qn("a:spcAft")):
        pPr.remove(old)
    sa = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(sa, qn("a:spcPts")).set("val", "0")
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    ls = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(ls, qn("a:spcPct")).set("val", "100000")


def _cell_write(cell, text: str, *,
                size: int = 7, bold: bool = False,
                color: RGBColor | None = None,
                align=PP_ALIGN.LEFT) -> None:
    """Write text into a table cell, clearing previous content."""
    _cell_margins(cell)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    _para_spacing_zero(p)
    for r in list(p.runs):
        p._p.remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _cell_write_multiline(cell, text: str, *,
                          size: int = 7, bold: bool = False,
                          color: RGBColor | None = None,
                          align=PP_ALIGN.LEFT) -> None:
    """Write multi-line text into a table cell."""
    _cell_margins(cell)
    tf = cell.text_frame
    tf.word_wrap = True
    for r in list(tf.paragraphs[0].runs):
        tf.paragraphs[0]._p.remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    lines = str(text).split("\n") if text else [""]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        _para_spacing_zero(p)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def _merge_cells(tbl, r1: int, c1: int, r2: int, c2: int) -> None:
    """Merge table cells from (r1,c1) to (r2,c2) inclusive."""
    tbl.cell(r1, c1).merge(tbl.cell(r2, c2))


def _set_table_borders(tbl, color_hex: str = "CCCCCC", size: int = 4) -> None:
    """Set thin borders on the whole table."""
    tblPr = tbl._tbl.tblPr
    if tblPr is None:
        tblPr = etree.SubElement(tbl._tbl, qn("a:tblPr"))
    borders_parts = [
        f'<a:top w:val="single" w:sz="{size}" w:space="0" w:color="{color_hex}"/>',
        f'<a:left w:val="single" w:sz="{size}" w:space="0" w:color="{color_hex}"/>',
        f'<a:bottom w:val="single" w:sz="{size}" w:space="0" w:color="{color_hex}"/>',
        f'<a:right w:val="single" w:sz="{size}" w:space="0" w:color="{color_hex}"/>',
        f'<a:insideH w:val="single" w:sz="{size}" w:space="0" w:color="{color_hex}"/>',
        f'<a:insideV w:val="single" w:sz="{size}" w:space="0" w:color="{color_hex}"/>',
    ]
    # python-pptx tables handle borders differently — set per-cell
    pass


def _add_textbox(slide, left, top, width, height, text,
                 font_size=10, bold=False, color=None, align=PP_ALIGN.LEFT,
                 font_name=FONT) -> None:
    """Add a textbox to a slide."""
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


# ─── Template loading ───────────────────────────────────────────────

def _load_base_presentation(template_path: str | Path | None = None) -> tuple:
    """Load the corporate template PPTX or create a blank."""
    tpl = Path(template_path) if template_path else TEMPLATE_PATH

    if tpl.exists():
        prs = Presentation(str(tpl))
        layout = prs.slides[0].slide_layout if len(prs.slides) > 0 else None
        # Delete existing slides
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            rId = sld_id.get(qn("r:id"))
            prs.part.drop_rel(rId)
            sld_id_lst.remove(sld_id)
        return prs, layout

    prs = Presentation()
    prs.slide_width = _W
    prs.slide_height = _H
    return prs, None


def _add_slide(prs: Presentation, layout) -> Any:
    """Add a slide using the template layout or blank."""
    if layout is not None:
        return prs.slides.add_slide(layout)
    blank_idx = min(6, len(prs.slide_layouts) - 1)
    return prs.slides.add_slide(prs.slide_layouts[blank_idx])


# ─── Stage 3: LLM slide spec ────────────────────────────────────────

_SLIDES_SYSTEM = """\
You are a Chief Procurement Officer designing a Supplier Impact Dashboard
presentation for a board-level audience. You are given JSON extracted from
the KB Input worksheet of a Supplier Impact Tracking file.
Generate a COMPLETE two-slide specification as JSON.

═══ SLIDE 1 — SUPPLIER IMPACT OVERVIEW ═══

This slide has FOUR sections:

1. EVALUATION SUMMARY — a single string like:
   "Evaluation: 48 possible suppliers being checked (15 already affected)"
   Count total suppliers in the data, and count those with severity R or Y as "affected".

2. COVERAGE DISTRIBUTION — aggregate supplier coverage data into 4 buckets:
   - no_coverage: count of suppliers with no/zero coverage (total_coverage_fg_days = 0 or null
     or "No coverage" or similar)
   - lt_4_days: count with coverage < 4 days
   - 5_to_15_days: count with coverage 5–15 days
   - gt_15_days: count with coverage > 15 days
   Use the total_coverage_fg_days or al_other_rm_coverage_days field.
   When coverage data is textual (e.g. "10 days", "2 weeks"), interpret it as days.
   When coverage is missing or unclear, count as no_coverage.

3. AFFECTED SUPPLIERS TABLE — one row per affected supplier (severity R or Y):
   Fields: supplier_name, cat, q_pave (int), l_pave (int), remarks
   - q_pave and l_pave: numeric risk scores (0–3). Use 0 if not available.
   - remarks: process impact summary, ≤ 120 chars.

4. ACTIONS TABLE — derive 4–8 governance actions from the data:
   Fields: action, resp, deadline, status_comments
   Focus on:
   - Follow-up on missing supplier feedback
   - Internal coverage assessment (ICO/DOM)
   - Alternate sourcing evaluation
   - Customer coverage & notification
   - Legal/FM notification status
   - Regulatory monitoring
   Keep actions concise and board-relevant.

5. CONTEXTUAL NOTES — 1–2 sentences explaining any caveats, e.g.:
   "15 of 48 suppliers confirmed as impacted. Sub-supplier dependencies
   remain under evaluation."

═══ SLIDE 2 — SUPPLIER FULFILLMENT DETAIL ═══

One row per impacted supplier with columns:
  supplier_name, host, material_planner, sda, coverage_date, coverage_after_actions,
  affected_product, customer, remarks

Content rules:
- coverage_date: use the date field or coverage date, format DD.MM.YYYY
- coverage_after_actions: qualitative (e.g. "No visibility", "ICO available",
  "Backup identified", "5 days additional")
- affected_product: from device_product_line field
- customer: combine dom_customer_name and ico_customer_name
- remarks: full root cause / mitigation details. Preserve ALL text — do NOT truncate.

═══ OUTPUT FORMAT ═══

Respond with ONLY valid JSON — no markdown fences, no explanation text.

CRITICAL JSON RULES:
  • No literal newlines, tabs inside string values — use "; " instead.
  • No trailing commas.
  • Must pass Python json.loads.

{{
  "presentation_title": "Supplier Situation Update",
  "last_update": "25.03.2026",
  "evaluation_summary": "Evaluation: 48 possible suppliers being checked (15 already affected)",
  "coverage_distribution": {{
    "no_coverage": 1,
    "lt_4_days": 10,
    "5_to_15_days": 13,
    "gt_15_days": 20
  }},
  "affected_suppliers": [
    {{
      "supplier_name": "Supplier Corp.",
      "cat": "SC",
      "q_pave": 2,
      "l_pave": 1,
      "remarks": "Casting process affected via sub-supplier dependency"
    }}
  ],
  "actions": [
    {{
      "action": "Closure on missing supplier feedbacks",
      "resp": "SDA team",
      "deadline": "ongoing",
      "status_comments": "Continuous follow-up with 3 remaining suppliers"
    }}
  ],
  "contextual_notes": "15 of 48 suppliers confirmed impacted. Sub-supplier dependencies under evaluation.",
  "supplier_details": [
    {{
      "supplier_name": "Supplier Corp.",
      "host": "Person A",
      "material_planner": "Person B",
      "sda": "Person C",
      "coverage_date": "25.03.2026",
      "coverage_after_actions": "ICO available until CW18",
      "affected_product": "Compressor, Actuator",
      "customer": "BMW, Mercedes",
      "remarks": "Root cause: gas supply disruption; dual sourcing under evaluation; interim air freight"
    }}
  ],
  "overall_severity": "HIGH",
  "footer_text": "T/SC SDA"
}}

IMPORTANT:
  • Include ALL affected suppliers in both tables.
  • overall_severity = "CRITICAL" if any severity R, "HIGH" if mix R+Y,
    "MEDIUM" if all Y, "LOW" if all G.
{glossary_block}"""

_SLIDES_USER = """\
Today: {today}
Current calendar week: {current_cw}

Generate the complete two-slide specification from this extracted data.

{extracted_json}"""


async def llm_generate_slide_spec(
    extracted: dict,
    llm_config: dict,
    session_id: str,
    glossary_context: str = "",
) -> dict:
    """Stage 3: LLM generates a complete slide specification."""
    llm = _create_llm(llm_config, max_tokens=16384)
    t0 = time.time()
    today = date.today().strftime("%d.%m.%Y")
    current_cw = extracted.get("current_cw", "CW??/????")

    n_input_suppliers = len(extracted.get("suppliers", []))
    print(f"\n[SID DEBUG] ═══ Stage 3: LLM slide spec ═══")
    print(f"[SID DEBUG] Input: {n_input_suppliers} suppliers, CW={current_cw}")

    glossary_block = (
        f"\n\nCOMPANY GLOSSARY:\n{glossary_context}" if glossary_context else ""
    )

    clean = {k: v for k, v in extracted.items() if not k.startswith("_")}
    extracted_json = json.dumps(clean, indent=2, ensure_ascii=False)

    messages = [
        SystemMessage(content=_SLIDES_SYSTEM.format(
            glossary_block=glossary_block,
        )),
        HumanMessage(content=_SLIDES_USER.format(
            today=today,
            current_cw=current_cw,
            extracted_json=extracted_json,
        )),
    ]

    try:
        response = await llm.ainvoke(messages)
        raw = response.content.strip()

        # ── Dump raw LLM response for debugging ──────────────────
        try:
            raw_dump_path = Path(__file__).parent / "uploads" / f"{session_id}_sid_slidespec_raw.txt"
            raw_dump_path.write_text(raw, encoding="utf-8")
            print(f"[SID DEBUG] Raw slide-spec LLM response saved: {raw_dump_path} "
                  f"({len(raw):,} chars)")
        except Exception as dump_exc:
            print(f"[SID DEBUG] Could not save raw slide-spec response: {dump_exc}")

        spec: dict = _parse_llm_json(raw, session_id=session_id)

        usage = response.response_metadata.get("token_usage", {})
        log_tokens(session_id, "sid_llm_slide_spec", usage,
                   llm_config.get("azure_deployment", ""))

        duration = (time.time() - t0) * 1000
        n_affected = len(spec.get("affected_suppliers", []))
        n_details = len(spec.get("supplier_details", []))
        n_actions = len(spec.get("actions", []))

        print(f"[SID DEBUG] Slide spec generated: {n_affected} affected suppliers, "
              f"{n_details} detail rows, {n_actions} actions "
              f"(severity={spec.get('overall_severity', '?')}, {duration:.0f}ms)")

        log_trace(session_id, "sid_llm_slide_spec",
                  f"Input: {len(extracted.get('suppliers', []))} suppliers",
                  f"Spec: {n_affected} affected, {n_details} detail rows, "
                  f"{n_actions} actions ({duration:.0f}ms)",
                  duration)
        return spec

    except json.JSONDecodeError as exc:
        duration = (time.time() - t0) * 1000
        print(f"[SID ERROR] Slide spec JSON parse failed — using deterministic fallback: "
              f"{str(exc)[:100]}")
        log_trace(session_id, "sid_llm_slide_spec",
                  "Generating slide spec",
                  f"JSON PARSE ERROR — using fallback: {str(exc)[:100]}",
                  duration, {"error": True, "fallback": True})
        return _deterministic_fallback_spec(extracted)

    except Exception as exc:
        duration = (time.time() - t0) * 1000
        print(f"[SID ERROR] Slide spec generation failed — using deterministic fallback: "
              f"{str(exc)[:100]}")
        log_trace(session_id, "sid_llm_slide_spec",
                  "Generating slide spec",
                  f"ERROR — using fallback: {str(exc)[:100]}",
                  duration, {"error": True, "fallback": True})
        return _deterministic_fallback_spec(extracted)


# ─── Deterministic fallback spec ────────────────────────────────────

def _deterministic_fallback_spec(extracted: dict) -> dict:
    """Compute slide spec without LLM from raw supplier data."""
    suppliers = extracted.get("suppliers", [])
    today = date.today().strftime("%d.%m.%Y")

    total = len(suppliers)
    affected = [s for s in suppliers
                if str(s.get("severity", "")).upper() in ("R", "Y")]
    n_affected = len(affected)

    # Coverage distribution
    no_cov = lt4 = s5_15 = gt15 = 0
    for s in suppliers:
        cov = s.get("total_coverage_fg_days") or s.get("al_other_rm_coverage_days")
        if cov is None or cov == "" or str(cov).lower() in ("nan", "none", "null", "n/a"):
            no_cov += 1
            continue
        try:
            days = float(str(cov).replace("days", "").replace("day", "").strip())
        except (ValueError, TypeError):
            no_cov += 1
            continue
        if days <= 0:
            no_cov += 1
        elif days < 4:
            lt4 += 1
        elif days <= 15:
            s5_15 += 1
        else:
            gt15 += 1

    # Build affected suppliers table
    affected_table = []
    for s in affected:
        affected_table.append({
            "supplier_name": s.get("vendor", ""),
            "cat": s.get("cat", ""),
            "q_pave": 0,
            "l_pave": 0,
            "remarks": (s.get("process_impacted", "") or "")[:120],
        })

    # Build details table
    details = []
    for s in affected:
        details.append({
            "supplier_name": s.get("vendor", ""),
            "host": "",
            "material_planner": s.get("category_buyer", ""),
            "sda": "",
            "coverage_date": s.get("date", ""),
            "coverage_after_actions": str(s.get("current_fuel_coverage", "") or ""),
            "affected_product": s.get("device_product_line", ""),
            "customer": "; ".join(filter(None, [
                s.get("dom_customer_name", ""),
                s.get("ico_customer_name", ""),
            ])),
            "remarks": s.get("remarks", ""),
        })

    # Default actions
    actions = [
        {"action": "Closure on missing supplier feedbacks", "resp": "SDA",
         "deadline": "ongoing", "status_comments": "Continuous follow-up"},
        {"action": "Internal coverage assessment (ICO)", "resp": "Planning",
         "deadline": "ongoing", "status_comments": "Evaluating current stock"},
        {"action": "Alternate sources evaluation", "resp": "Purchasing",
         "deadline": "ongoing", "status_comments": "Global check in progress"},
        {"action": "Customer notification", "resp": "Sales",
         "deadline": "TBD", "status_comments": "Pending FM letters"},
    ]

    severities = [str(s.get("severity", "")).upper() for s in suppliers]
    if "R" in severities:
        overall = "HIGH"
    elif "Y" in severities:
        overall = "MEDIUM"
    else:
        overall = "LOW"

    return {
        "presentation_title": "Supplier Situation Update",
        "last_update": today,
        "evaluation_summary": f"Evaluation: {total} possible suppliers being checked ({n_affected} already affected)",
        "coverage_distribution": {
            "no_coverage": no_cov,
            "lt_4_days": lt4,
            "5_to_15_days": s5_15,
            "gt_15_days": gt15,
        },
        "affected_suppliers": affected_table,
        "actions": actions,
        "contextual_notes": f"{n_affected} of {total} suppliers confirmed impacted.",
        "supplier_details": details,
        "overall_severity": overall,
        "footer_text": "T/SC SDA",
        "_fallback": True,
    }


# ═══════════════════════════════════════════════════════════════════════
#  Stage 4: PPTX renderer
# ═══════════════════════════════════════════════════════════════════════

# ─── Slide 1 renderer ───────────────────────────────────────────────

def _render_slide1(prs, layout, spec: dict) -> None:
    """Render the executive overview slide."""
    slide = _add_slide(prs, layout)

    last_update = spec.get("last_update", date.today().strftime("%d.%m.%Y"))
    eval_summary = spec.get("evaluation_summary", "")
    cov = spec.get("coverage_distribution", {})
    affected = spec.get("affected_suppliers", [])
    actions = spec.get("actions", [])
    notes = spec.get("contextual_notes", "")
    footer_text = spec.get("footer_text", "T/SC SDA")

    # ── Clear any existing placeholders ──
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # ── Title ──
    _add_textbox(slide,
                 Emu(2475127), Emu(545661), Emu(2031266), Emu(152400),
                 "Detail of the affected suppliers",
                 font_size=14, bold=False, color=_C["DARK_BLUE"])

    # ── Evaluation summary (top bar) ──
    _add_textbox(slide,
                 Inches(0.26), Inches(0.17), Inches(9.47), Inches(0.32),
                 eval_summary,
                 font_size=11, bold=True, color=_C["DARK_BLUE"])

    # ── Last Update (top right) ──
    _add_textbox(slide,
                 Emu(7530340), Emu(418000), Emu(1463040), Emu(150000),
                 f"Last Update: {last_update}",
                 font_size=9, bold=True, color=_C["RED"],
                 align=PP_ALIGN.RIGHT)

    # ── "Overview of suppliers" sub-title ──
    _add_textbox(slide,
                 Inches(0.43), Inches(0.63), Inches(1.6), Inches(0.15),
                 "Overview of suppliers",
                 font_size=9, bold=True, color=_C["DARK_BLUE"])

    # ── Coverage legend (left panel, below chart area) ──
    legend_items = [
        ("No Coverage", _C["BLACK"]),
        ("Coverage < 4 days", _C["RED"]),
        ("Coverage 5 to 15 days", _C["YELLOW"]),
        ("Coverage > 15 days", _C["GREEN"]),
    ]
    cov_values = [
        cov.get("no_coverage", 0),
        cov.get("lt_4_days", 0),
        cov.get("5_to_15_days", 0),
        cov.get("gt_15_days", 0),
    ]

    # Draw pie chart — one category per coverage bucket
    chart_data = CategoryChartData()
    chart_data.categories = ["No Coverage", "< 4 days", "5–15 days", "> 15 days"]
    # Ensure all values > 0 for pie chart visibility
    chart_vals = [max(v, 0.001) for v in cov_values]
    chart_data.add_series("Coverage", chart_vals)

    chart_left = Inches(0.50)
    chart_top = Inches(0.81)
    chart_w = Inches(1.63)
    chart_h = Inches(1.63)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, chart_left, chart_top, chart_w, chart_h,
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = False

    # Set pie slice colors
    plot = chart.plots[0]
    series = plot.series[0]
    pie_colors = [_C["BLACK"], _C["RED"], _C["YELLOW"], _C["GREEN"]]
    for i, clr in enumerate(pie_colors):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = clr

    # Coverage count labels around chart
    count_positions = [
        (Inches(1.21), Inches(0.92)),   # no_coverage count
        (Inches(1.66), Inches(1.14)),   # lt_4 count
        (Inches(1.33), Inches(2.14)),   # 5_15 count
        (Inches(0.68), Inches(1.26)),   # gt_15 count
    ]
    for idx, (cx, cy) in enumerate(count_positions):
        _add_textbox(slide, cx, cy, Inches(0.19), Inches(0.17),
                     str(cov_values[idx]),
                     font_size=8, bold=True, color=_C["WHITE"])

    # Legend items
    legend_top_start = Inches(2.68)
    for i, (label, clr) in enumerate(legend_items):
        y = legend_top_start + Inches(i * 0.156)
        # Colored rectangle
        rect = slide.shapes.add_shape(
            1, Inches(0.41), y, Inches(0.12), Inches(0.09)  # 1 = RECTANGLE
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = clr
        rect.line.fill.background()
        # Label text
        _add_textbox(slide, Inches(0.58), y, Inches(0.85), Inches(0.10),
                     f"{label}", font_size=6, bold=False)

    # ── Affected Suppliers Table (right panel) ──
    n_sup_rows = max(len(affected), 1)
    sup_tbl_rows = n_sup_rows + 1  # +1 for header
    sup_col_widths = [Emu(1612645), Emu(543996), Emu(492369), Emu(614859), Emu(3104543)]

    sup_tbl_shape = slide.shapes.add_table(
        sup_tbl_rows, 5,
        Emu(2429300), Emu(683000),
        int(sum(sup_col_widths)), Emu(min(sup_tbl_rows * 130000 + 143107, 2300000)),
    )
    sup_tbl = sup_tbl_shape.table

    # Disable banding
    tblPr = sup_tbl._tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("lastRow", "0")

    for ci, w in enumerate(sup_col_widths):
        sup_tbl.columns[ci].width = w

    # Header
    sup_headers = ["Supplier name", "Cat", "Q-PAVE", "L-\nPAVE",
                    "Remarks for affected process & partner"]
    for ci, label in enumerate(sup_headers):
        cell = sup_tbl.cell(0, ci)
        _cell_fill(cell, _C["HDR_BLUE"])
        _cell_write_multiline(cell, label, size=7, bold=True, color=_C["WHITE"],
                              align=PP_ALIGN.CENTER if ci in (1, 2, 3) else PP_ALIGN.LEFT)
    sup_tbl.rows[0].height = Emu(143107)

    # Data rows
    for ri, sup in enumerate(affected):
        row_idx = ri + 1
        _cell_write(sup_tbl.cell(row_idx, 0), sup.get("supplier_name", ""), size=6)
        _cell_write(sup_tbl.cell(row_idx, 1), sup.get("cat", ""), size=6,
                    align=PP_ALIGN.CENTER)
        _cell_write(sup_tbl.cell(row_idx, 2), str(sup.get("q_pave", "")), size=6,
                    align=PP_ALIGN.CENTER)
        _cell_write(sup_tbl.cell(row_idx, 3), str(sup.get("l_pave", "")), size=6,
                    align=PP_ALIGN.CENTER)
        _cell_write(sup_tbl.cell(row_idx, 4), sup.get("remarks", ""), size=6)
        sup_tbl.rows[row_idx].height = Emu(120000)

    # ── "Actions moving forward" title ──
    _add_textbox(slide,
                 Inches(0.30), Inches(3.42), Inches(1.90), Inches(0.17),
                 "Actions moving forward",
                 font_size=14, bold=False, color=_C["DARK_BLUE"])

    # ── Actions Table ──
    n_act_rows = max(len(actions), 1)
    act_tbl_rows = n_act_rows + 1
    act_col_widths = [Emu(1613660), Emu(1029335), Emu(771525), Emu(5138073)]

    act_tbl_top = Inches(3.56)
    act_tbl_shape = slide.shapes.add_table(
        act_tbl_rows, 4,
        Emu(284163), int(act_tbl_top),
        int(sum(act_col_widths)), Emu(min(act_tbl_rows * 100000 + 95149, 1600000)),
    )
    act_tbl = act_tbl_shape.table

    tblPr2 = act_tbl._tbl.tblPr
    tblPr2.set("bandRow", "0")
    tblPr2.set("bandCol", "0")
    tblPr2.set("firstRow", "0")
    tblPr2.set("lastRow", "0")

    for ci, w in enumerate(act_col_widths):
        act_tbl.columns[ci].width = w

    # Header
    act_headers = ["Action", "Resp.", "Deadline", "Status / Comments"]
    for ci, label in enumerate(act_headers):
        cell = act_tbl.cell(0, ci)
        _cell_fill(cell, _C["HDR_BLUE"])
        _cell_write(cell, label, size=7, bold=True, color=_C["WHITE"])
    act_tbl.rows[0].height = Emu(95149)

    # Data rows
    for ri, act in enumerate(actions):
        row_idx = ri + 1
        _cell_write(act_tbl.cell(row_idx, 0), act.get("action", ""), size=6)
        _cell_write(act_tbl.cell(row_idx, 1), act.get("resp", ""), size=6)
        _cell_write(act_tbl.cell(row_idx, 2), act.get("deadline", ""), size=6)
        _cell_write(act_tbl.cell(row_idx, 3), act.get("status_comments", ""), size=6)
        act_tbl.rows[row_idx].height = Emu(100000)

    # ── Footer ──
    month_str = datetime.now().strftime("%m/%Y")
    _add_textbox(slide,
                 Inches(5.08), Inches(5.39), Inches(4.30), Inches(0.16),
                 f"{footer_text} – {month_str}",
                 font_size=7, color=_C["GREY"])

    # Page number
    _add_textbox(slide,
                 Inches(9.41), Inches(5.39), Inches(0.32), Inches(0.16),
                 "1", font_size=7, color=_C["GREY"], align=PP_ALIGN.RIGHT)


# ─── Slide 2 renderer ───────────────────────────────────────────────

def _render_slide2(prs, layout, spec: dict) -> None:
    """Render the supplier detail slide."""
    slide = _add_slide(prs, layout)

    details = spec.get("supplier_details", [])
    last_update = spec.get("last_update", date.today().strftime("%d.%m.%Y"))
    footer_text = spec.get("footer_text", "T/SC SDA")

    # Clear placeholders
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)

    # ── Title ──
    _add_textbox(slide,
                 Inches(0.09), Inches(0.08), Inches(8.10), Inches(0.23),
                 "Fulfillment Overview (impacted suppliers)",
                 font_size=14, bold=True, color=_C["DARK_BLUE"])

    # ── Last Update ──
    _add_textbox(slide,
                 Emu(7591706), Emu(67564), Emu(1463040), Emu(150000),
                 f"Last Update: {last_update}",
                 font_size=9, bold=True, color=_C["RED"],
                 align=PP_ALIGN.RIGHT)

    # ── Accent bar ──
    bar = slide.shapes.add_shape(
        1, Inches(0.09), Inches(0.25), Inches(3.43), Emu(52917)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _C["HDR_BLUE"]
    bar.line.fill.background()

    # ── Supplier Detail Table ──
    n_rows = max(len(details), 1)
    total_rows = n_rows + 1  # +1 header
    col_widths_emu = [
        Emu(792480),   # Supplier name
        Emu(678180),   # Host
        Emu(640080),   # Material Planner
        Emu(558165),   # SDA
        Emu(857250),   # Coverage Date
        Emu(932688),   # Coverage after actions
        Emu(791845),   # Affected product
        Emu(880110),   # Customer
        Emu(2841938),  # Remarks
    ]

    tbl_left = Emu(84138)
    tbl_top = Emu(278765)
    tbl_w = sum(col_widths_emu)

    # Dynamic row heights
    header_h = Emu(380930)
    available_h = _H - tbl_top - Inches(0.25)
    data_row_h = int((available_h - header_h) / max(n_rows, 1))
    data_row_h = max(data_row_h, Emu(200000))
    data_row_h = min(data_row_h, Emu(660000))

    tbl_h = header_h + data_row_h * n_rows

    tbl_shape = slide.shapes.add_table(
        total_rows, 9,
        int(tbl_left), int(tbl_top), int(tbl_w), int(tbl_h),
    )
    tbl = tbl_shape.table

    tblPr = tbl._tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("lastRow", "0")

    for ci, w in enumerate(col_widths_emu):
        tbl.columns[ci].width = w

    tbl.rows[0].height = int(header_h)

    # Header
    detail_headers = [
        "Supplier\nname", "Host", "Material\nPlanner", "SDA",
        "Coverage\nDate", "Coverage\nafter actions",
        "Affected\nproduct", "Customer", "Remarks",
    ]
    for ci, label in enumerate(detail_headers):
        cell = tbl.cell(0, ci)
        _cell_fill(cell, _C["DARK_BLUE"])
        _cell_write_multiline(cell, label, size=8, bold=True, color=_C["WHITE"],
                              align=PP_ALIGN.CENTER)

    # Data rows
    for ri, row in enumerate(details):
        row_idx = ri + 1
        tbl.rows[row_idx].height = data_row_h

        fields = [
            ("supplier_name", PP_ALIGN.LEFT),
            ("host", PP_ALIGN.LEFT),
            ("material_planner", PP_ALIGN.LEFT),
            ("sda", PP_ALIGN.LEFT),
            ("coverage_date", PP_ALIGN.CENTER),
            ("coverage_after_actions", PP_ALIGN.LEFT),
            ("affected_product", PP_ALIGN.LEFT),
            ("customer", PP_ALIGN.LEFT),
            ("remarks", PP_ALIGN.LEFT),
        ]
        for ci, (field, al) in enumerate(fields):
            text = row.get(field, "") or ""
            _cell_write_multiline(tbl.cell(row_idx, ci), str(text),
                                  size=6, align=al)

        # Alternate row shading
        if ri % 2 == 1:
            for ci in range(9):
                _cell_fill(tbl.cell(row_idx, ci), _C["LIGHT_GREY"])

    # ── Footer ──
    month_str = datetime.now().strftime("%m/%Y")
    _add_textbox(slide,
                 Inches(5.08), Inches(5.39), Inches(4.30), Inches(0.16),
                 f"{footer_text} – {month_str}",
                 font_size=7, color=_C["GREY"])

    _add_textbox(slide,
                 Inches(9.41), Inches(5.39), Inches(0.32), Inches(0.16),
                 "2", font_size=7, color=_C["GREY"], align=PP_ALIGN.RIGHT)


# ─── Main renderer ──────────────────────────────────────────────────

def render_pptx_from_spec(slide_spec: dict,
                          output_path: str | None = None,
                          template_path: str | Path | None = None,
                          ) -> io.BytesIO:
    """
    Stage 4 (deterministic renderer): Convert slide spec to PPTX.
    Always produces exactly 2 slides.
    """
    prs, layout = _load_base_presentation(template_path)

    _render_slide1(prs, layout, slide_spec)
    _render_slide2(prs, layout, slide_spec)

    slide_spec["slide_count"] = 2

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    if output_path:
        with open(output_path, "wb") as fh:
            fh.write(buf.read())
        buf.seek(0)

    return buf


# ─── Combined entry point ───────────────────────────────────────────

async def generate_sid_dashboard(
    extracted_data: dict,
    llm_config: dict,
    session_id: str,
    output_path: str | None = None,
    glossary_context: str = "",
) -> tuple[io.BytesIO, dict]:
    """
    Full SID dashboard generation pipeline.

    1. LLM generates slide specification.
    2. PPTX renderer converts spec to file.

    Returns (buf: BytesIO, slide_spec: dict)
    """
    spec = await llm_generate_slide_spec(
        extracted_data, llm_config, session_id, glossary_context
    )
    buf = render_pptx_from_spec(spec, output_path=output_path)
    return buf, spec
