"""
PPT Parser Module
=================
Methodical, section-by-section extraction of PowerPoint content.
Outputs structured JSON for each slide and section.
"""

import json
import re
from pathlib import Path
from typing import Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_CHART_TYPE
from pptx.table import Table


def _extract_table(table: Table) -> list[list[str]]:
    """Extract table data as a 2D list of strings."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)
    return rows


def _table_to_markdown(table_data: list[list[str]]) -> str:
    """Convert 2D table data to markdown table string."""
    if not table_data or not table_data[0]:
        return ""
    header = table_data[0]
    md = "| " + " | ".join(header) + " |\n"
    md += "| " + " | ".join(["---"] * len(header)) + " |\n"
    for row in table_data[1:]:
        # Pad row if shorter than header
        padded = row + [""] * (len(header) - len(row))
        md += "| " + " | ".join(padded[:len(header)]) + " |\n"
    return md


def _extract_chart(chart) -> dict[str, Any]:
    """Extract chart data including series, categories, and values."""
    chart_data = {
        "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
        "title": "",
        "series": []
    }
    try:
        if chart.has_title and chart.chart_title and chart.chart_title.has_text_frame:
            chart_data["title"] = chart.chart_title.text_frame.text
    except Exception:
        pass

    try:
        plot = chart.plots[0]
        categories = []
        try:
            categories = [str(c) for c in plot.categories]
        except Exception:
            pass
        chart_data["categories"] = categories

        for series in plot.series:
            series_info = {"name": "", "values": []}
            try:
                if hasattr(series, 'tx') and series.tx and hasattr(series.tx, 'strRef'):
                    pass  # Complex extraction
            except Exception:
                pass
            try:
                series_info["values"] = [
                    round(float(v), 2) if v is not None else None
                    for v in series.values
                ]
            except Exception:
                pass
            chart_data["series"].append(series_info)
    except Exception as e:
        chart_data["extraction_error"] = str(e)

    return chart_data


def _extract_text_from_shape(shape) -> str:
    """Extract all text from a shape's text frame."""
    if not shape.has_text_frame:
        return ""
    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            # Preserve bullet level info
            level = para.level if para.level else 0
            prefix = "  " * level + ("• " if level > 0 else "")
            paragraphs.append(prefix + text)
    return "\n".join(paragraphs)


def _detect_color_coding(shape) -> list[dict]:
    """Detect RAG (Red/Amber/Green) color coding in table cells."""
    color_info = []
    if not shape.has_table:
        return color_info
    for row_idx, row in enumerate(shape.table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                fill = cell.fill
                if fill and fill.type is not None:
                    fore_color = None
                    try:
                        fore_color = str(fill.fore_color.rgb) if fill.fore_color and fill.fore_color.rgb else None
                    except Exception:
                        pass
                    if fore_color:
                        # Map common RAG colors
                        color_name = _map_rgb_to_rag(fore_color)
                        if color_name:
                            color_info.append({
                                "row": row_idx,
                                "col": col_idx,
                                "cell_text": cell.text.strip(),
                                "color": color_name,
                                "rgb": fore_color
                            })
            except Exception:
                continue
    return color_info


def _map_rgb_to_rag(rgb: str) -> str | None:
    """Map RGB hex to RAG status color name."""
    rgb_upper = rgb.upper()
    r = int(rgb_upper[0:2], 16)
    g = int(rgb_upper[2:4], 16)
    b = int(rgb_upper[4:6], 16)

    # Red-ish
    if r > 180 and g < 100 and b < 100:
        return "RED"
    # Green-ish
    if g > 150 and r < 150 and b < 100:
        return "GREEN"
    # Amber/Yellow-ish
    if r > 180 and g > 150 and b < 100:
        return "AMBER"
    # Orange
    if r > 200 and 80 < g < 170 and b < 80:
        return "AMBER"
    return None


def extract_slide(slide, slide_index: int) -> dict[str, Any]:
    """Extract all content from a single slide."""
    slide_data = {
        "slide_number": slide_index + 1,
        "title": "",
        "subtitle": "",
        "text_content": [],
        "tables": [],
        "charts": [],
        "color_coding": [],
        "notes": ""
    }

    # Extract slide title
    if slide.shapes.title:
        slide_data["title"] = slide.shapes.title.text.strip()

    # Extract notes
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()

    for shape in slide.shapes:
        # Skip title shape (already extracted)
        if shape == slide.shapes.title:
            continue

        # Tables
        if shape.has_table:
            table_data = _extract_table(shape.table)
            if table_data:
                slide_data["tables"].append({
                    "data": table_data,
                    "markdown": _table_to_markdown(table_data),
                    "row_count": len(table_data),
                    "col_count": len(table_data[0]) if table_data else 0
                })
                # Check for color coding
                colors = _detect_color_coding(shape)
                if colors:
                    slide_data["color_coding"].extend(colors)

        # Charts
        elif shape.has_chart:
            chart_data = _extract_chart(shape.chart)
            slide_data["charts"].append(chart_data)

        # Text
        elif shape.has_text_frame:
            text = _extract_text_from_shape(shape)
            if text:
                # Check if this might be a subtitle
                if not slide_data["subtitle"] and shape != slide.shapes.title:
                    # Heuristic: short text near top might be subtitle
                    if shape.top and shape.top < Emu(2000000) and len(text) < 100:
                        slide_data["subtitle"] = text
                        continue
                slide_data["text_content"].append(text)

    return slide_data


def detect_sections(slides_data: list[dict]) -> list[dict]:
    """
    Detect logical sections in the presentation.
    Uses the Agenda slide and section header patterns.
    """
    sections = []
    agenda_items = []
    current_section = None

    # First pass: find agenda slide
    for slide in slides_data:
        title_lower = slide["title"].lower()
        if any(kw in title_lower for kw in ["agenda", "table of contents", "contents", "overview"]):
            # Extract agenda items from text content
            for text in slide["text_content"]:
                for line in text.split("\n"):
                    line = line.strip().lstrip("•–-·→► 0123456789.")
                    if line and len(line) > 2:
                        agenda_items.append(line.strip())
            break

    # Section detection keywords (maps to known section types)
    section_keywords = {
        "crisis": "Ongoing Crisis/Issue Summary",
        "issue summary": "Ongoing Crisis/Issue Summary",
        "current situation": "Ongoing Crisis/Issue Summary",
        "supplier": "Supplier Situation",
        "production": "Production Situation",
        "customer fulfilment": "Global Customer Fulfilment",
        "customer fulfillment": "Global Customer Fulfilment",
        "fulfilment dashboard": "Global Customer Fulfilment",
        "fulfillment dashboard": "Global Customer Fulfilment",
        "supply coverage": "Global Customer Fulfilment",
        "customer situation": "Customer Situation & Demand",
        "demand development": "Customer Situation & Demand",
        "demand": "Customer Situation & Demand",
        "freight": "Freight & Logistics",
        "logistics": "Freight & Logistics",
        "capacity situation": "Freight & Logistics",
        "cost impact": "Cost Impact",
        "cost": "Cost Impact",
        "financial": "Cost Impact",
        "appendix": "Appendix",
        "backup": "Appendix",
    }

    # Second pass: assign slides to sections
    for slide in slides_data:
        title_lower = slide["title"].lower()

        # Skip title slide and agenda
        if any(kw in title_lower for kw in ["agenda", "table of contents"]):
            if not current_section:
                current_section = {
                    "section_name": "Introduction / Agenda",
                    "slides": [slide]
                }
                sections.append(current_section)
            continue

        # Check if this slide starts a new section
        matched_section = None
        for keyword, section_name in section_keywords.items():
            if keyword in title_lower:
                matched_section = section_name
                break

        # Also match against agenda items
        if not matched_section:
            for agenda_item in agenda_items:
                if (agenda_item.lower() in title_lower or
                    title_lower in agenda_item.lower() or
                    any(w in title_lower for w in agenda_item.lower().split() if len(w) > 4)):
                    matched_section = agenda_item
                    break

        if matched_section:
            current_section = {
                "section_name": matched_section,
                "slides": [slide]
            }
            sections.append(current_section)
        elif current_section:
            current_section["slides"].append(slide)
        else:
            # No section yet, create a general one
            current_section = {
                "section_name": slide["title"] if slide["title"] else "Introduction",
                "slides": [slide]
            }
            sections.append(current_section)

    # If no sections were detected, group all slides into one
    if not sections:
        sections = [{
            "section_name": "Full Presentation",
            "slides": slides_data
        }]

    return sections


def format_section_for_llm(section: dict) -> str:
    """Format a section's content as a structured text block for the LLM."""
    output = []
    output.append(f"## SECTION: {section['section_name']}")
    output.append(f"Number of slides in section: {len(section['slides'])}")
    output.append("")

    for slide in section["slides"]:
        output.append(f"### Slide {slide['slide_number']}: {slide['title']}")
        if slide["subtitle"]:
            output.append(f"Subtitle: {slide['subtitle']}")

        if slide["text_content"]:
            output.append("**Text Content:**")
            for text in slide["text_content"]:
                output.append(text)

        if slide["tables"]:
            for i, table in enumerate(slide["tables"]):
                output.append(f"**Table {i+1}** ({table['row_count']}x{table['col_count']}):")
                output.append(table["markdown"])

        if slide["charts"]:
            for i, chart in enumerate(slide["charts"]):
                output.append(f"**Chart {i+1}:** {chart.get('title', 'Untitled')}")
                output.append(f"  Type: {chart.get('chart_type', 'unknown')}")
                if chart.get("categories"):
                    output.append(f"  Categories: {', '.join(chart['categories'][:20])}")
                for s_idx, series in enumerate(chart.get("series", [])):
                    vals = series.get("values", [])
                    output.append(f"  Series {s_idx+1} ({series.get('name', '')}): {vals[:20]}")

        if slide["color_coding"]:
            output.append("**RAG Status Indicators:**")
            for cc in slide["color_coding"]:
                output.append(f"  [{cc['color']}] Row {cc['row']}, Col {cc['col']}: {cc['cell_text']}")

        if slide["notes"]:
            output.append(f"**Speaker Notes:** {slide['notes']}")

        output.append("---")

    return "\n".join(output)


def parse_presentation(filepath: str) -> dict[str, Any]:
    """
    Main entry point: parse entire presentation and return structured data.
    Returns sections with formatted content ready for LLM processing.
    """
    prs = Presentation(filepath)
    total_slides = len(prs.slides)

    # Extract all slides
    slides_data = []
    for idx, slide in enumerate(prs.slides):
        slide_data = extract_slide(slide, idx)
        slides_data.append(slide_data)

    # Detect sections
    sections = detect_sections(slides_data)

    # Format each section
    result = {
        "filename": Path(filepath).name,
        "total_slides": total_slides,
        "total_sections": len(sections),
        "sections": []
    }

    for section in sections:
        formatted = format_section_for_llm(section)
        result["sections"].append({
            "section_name": section["section_name"],
            "slide_count": len(section["slides"]),
            "slide_numbers": [s["slide_number"] for s in section["slides"]],
            "formatted_content": formatted,
            "raw_slides": section["slides"]
        })

    return result
