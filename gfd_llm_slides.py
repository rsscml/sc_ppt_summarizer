"""
GFD LLM Slide Generator
=======================
Two-stage pipeline:

Stage 1 — LLM slide spec:
  Receives the structured JSON from gfd_llm_parser and generates a complete
  slide specification — RAG colors, condensed cell text, groupings, KB coverage
  strings. The LLM makes ALL content decisions.

Stage 2 — PPTX renderer:
  A thin "paint by numbers" renderer that converts the slide spec JSON into a
  python-pptx Presentation object that matches the corporate GFD template:
  - 10 × 5.625 inch slides (standard 16:9)
  - 14-column table: Product Group | Plant | Customer/Channel | KB Coverage |
    6 × CW | Quarter | Supplier | Comment | FM Detail Letter
  - 2-row merged header (dark-blue text headers + green CW band)
  - Product-group column with vertical merge and amber background
  - RAG-coloured CW cells (GREEN / AMBER / RED)
  - Corporate font (Arial) and colour scheme

If Stage 1 fails, a deterministic fallback computes RAG colors arithmetically
from the integer coverage CW fields produced by the extractor, so the user
always gets a working dashboard even on LLM error.

When a corporate template PPTX is available, it is loaded as the base
presentation so that theme colours, logos, and slide-master styling are
preserved exactly.
"""

from __future__ import annotations

import io
import json
import math
import time
from datetime import date, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

from langchain_openai import AzureChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage

from agent import log_tokens, log_trace


# ─── Template path ───────────────────────────────────────────────────

TEMPLATE_PATH = Path(__file__).parent / "gfd_ppt_template.pptx"


# ─── LLM factory ────────────────────────────────────────────────────

def _create_llm(config: dict, max_tokens: int = 64000) -> AzureChatOpenAI:
    return AzureChatOpenAI(
        azure_deployment=config["azure_deployment"],
        azure_endpoint=config["azure_endpoint"],
        api_key=config["api_key"],
        api_version=config.get("api_version", "2024-12-01-preview"),
        #temperature=0.15,
        max_tokens=max_tokens,
    )


# ─── Colour palette (matching the corporate GFD template) ───────────

_C: dict[str, RGBColor] = {
    "HDR_BLUE":     RGBColor(0x00, 0x45, 0x7E),   # dk2/tx2 — header text columns
    "HDR_GREEN":    RGBColor(0x00, 0x80, 0x00),   # CW header band
    "GREEN":        RGBColor(0x00, 0xB0, 0x50),   # RAG green (data cells)
    "AMBER":        RGBColor(0xFF, 0xC0, 0x00),   # RAG amber + product-group bg
    "RED":          RGBColor(0xFF, 0x00, 0x00),   # RAG red
    "GREY":         RGBColor(0xA0, 0xA0, 0xA0),   # no-data
    "WHITE":        RGBColor(0xFF, 0xFF, 0xFF),
    "BLACK":        RGBColor(0x00, 0x00, 0x00),
}

_RAG_BG: dict[str, RGBColor] = {
    "GREEN": _C["GREEN"],
    "AMBER": _C["AMBER"],
    "RED":   _C["RED"],
    "GREY":  _C["GREY"],
}

# Slide dimensions — standard 16:9 matching the template
_W  = Inches(10.0)
_H  = Inches(5.625)

# Table position and geometry — matched from template
_TBL_LEFT = Emu(200938)     # 0.220 in
_TBL_TOP  = Emu(388775)     # 0.425 in

# Title position — matched from template placeholder
_TITLE_LEFT   = Emu(200938)
_TITLE_TOP    = Emu(111776)
_TITLE_WIDTH  = Emu(8639908)
_TITLE_HEIGHT = Emu(276999)

FONT = "Arial"

# Column widths in EMU — matched from template
_COL_WIDTHS_EMU = [
    882650,    # col 0: Product Group (PG)     — 0.965 in
    312545,    # col 1: Plant                  — 0.342 in
    1524318,   # col 2: Customer / Channel     — 1.667 in
    576611,    # col 3: KB Coverage (CW/YY)    — 0.631 in
    192832,    # col 4: CW+0                   — 0.211 in
    192832,    # col 5: CW+1
    192832,    # col 6: CW+2
    192832,    # col 7: CW+3
    192832,    # col 8: CW+4
    192832,    # col 9: CW+5
    192832,    # col 10: Quarter               — 0.211 in
    1785894,   # col 11: Supplier              — 1.953 in
    1132964,   # col 12: Comment               — 1.239 in
    698178,    # col 13: FM Detail Letter       — 0.764 in
]
_N_COLS        = len(_COL_WIDTHS_EMU)
_N_CW_COLS     = 6          # CW columns (cols 4–9)
_CW_COL_START  = 4          # first CW column index
_Q_COL         = 10         # quarter column index
_N_FIXED_LEFT  = 4          # text columns before CW band (0–3)
_N_FIXED_RIGHT = 3          # text columns after CW band (11–13)

# Max data rows per slide (excluding 2-row header).  The template comfortably
# fits ~24 data rows within 4.8 in of table height.
_MAX_ROWS_PER_SLIDE = 24

# Header row height
_HDR_ROW_H = Emu(231700)    # row 0 — 0.253 in
_HDR_SUB_H = Emu(385020)    # row 1 — 0.421 in (CW numbers row)


# ─── Low-level python-pptx helpers ──────────────────────────────────

def _cell_fill(cell, rgb: RGBColor) -> None:
    """Set solid background fill on a table cell via XML."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for child in list(tcPr):
        local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if local in ("solidFill", "gradFill", "noFill", "pattFill", "blipFill"):
            tcPr.remove(child)
    sf = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(sf, qn("a:srgbClr"))
    clr.set("val", f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")


def _cell_margins_zero(cell, anchor: str = "ctr") -> None:
    """Set cell-level margins to zero and vertical anchor via tcPr XML.
    This matches the corporate template where all data cells have marL/R/T/B=0."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set("marL", "0")
    tcPr.set("marR", "0")
    tcPr.set("marT", "0")
    tcPr.set("marB", "0")
    tcPr.set("anchor", anchor)


def _para_spacing_zero(p) -> None:
    """Set paragraph spacing to 0pt before/after and 100% line spacing."""
    pPr = p._p.get_or_add_pPr()
    # spcBef
    for old in pPr.findall(qn("a:spcBef")):
        pPr.remove(old)
    sb = etree.SubElement(pPr, qn("a:spcBef"))
    etree.SubElement(sb, qn("a:spcPts")).set("val", "0")
    # spcAft
    for old in pPr.findall(qn("a:spcAft")):
        pPr.remove(old)
    sa = etree.SubElement(pPr, qn("a:spcAft"))
    etree.SubElement(sa, qn("a:spcPts")).set("val", "0")
    # lnSpc
    for old in pPr.findall(qn("a:lnSpc")):
        pPr.remove(old)
    ls = etree.SubElement(pPr, qn("a:lnSpc"))
    etree.SubElement(ls, qn("a:spcPct")).set("val", "100000")


def _cell_write(cell, text: str, *,
                size: int = 6, bold: bool = False,
                color: RGBColor | None = None,
                align=PP_ALIGN.LEFT) -> None:
    """Write text into a table cell, clearing previous content."""
    _cell_margins_zero(cell)
    tf = cell.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = align
    _para_spacing_zero(p)
    for r in list(p.runs):
        p._p.remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _cell_write_multiline(cell, text: str, *,
                          size: int = 6, bold: bool = False,
                          color: RGBColor | None = None,
                          align=PP_ALIGN.LEFT) -> None:
    """Write multi-line text (newline-separated) into a table cell."""
    _cell_margins_zero(cell)
    tf = cell.text_frame
    tf.word_wrap = True

    # Clear existing content
    for r in list(tf.paragraphs[0].runs):
        tf.paragraphs[0]._p.remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)

    lines = str(text).split("\n") if text else [""]
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        _para_spacing_zero(p)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color


def _merge_cells(tbl, r1: int, c1: int, r2: int, c2: int) -> None:
    """Merge table cells from (r1,c1) to (r2,c2) inclusive."""
    tbl.cell(r1, c1).merge(tbl.cell(r2, c2))


# ─── Template loading ───────────────────────────────────────────────

def _load_base_presentation(template_path: str | Path | None = None) -> tuple:
    """
    Load the corporate template PPTX as the base presentation.

    Returns (prs, layout) where layout is the slide layout to use for
    new slides.  If the template is unavailable, creates a blank
    presentation with matching dimensions and returns (prs, None).
    """
    tpl = Path(template_path) if template_path else TEMPLATE_PATH

    if tpl.exists():
        prs = Presentation(str(tpl))

        # Capture the layout used by the first (example) slide
        layout = prs.slides[0].slide_layout if len(prs.slides) > 0 else None

        # Delete all existing slides (they contain example data)
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            rId = sld_id.get(qn("r:id"))
            prs.part.drop_rel(rId)
            sld_id_lst.remove(sld_id)

        return prs, layout

    # Fallback: create a blank presentation with matching dimensions
    prs = Presentation()
    prs.slide_width  = _W
    prs.slide_height = _H
    return prs, None


def _add_slide(prs: Presentation, layout) -> Any:
    """Add a slide using the template layout (or blank if unavailable)."""
    if layout is not None:
        return prs.slides.add_slide(layout)
    # Fallback: use blank layout (index 6 is typically blank)
    blank_idx = min(6, len(prs.slide_layouts) - 1)
    return prs.slides.add_slide(prs.slide_layouts[blank_idx])


# ─── Stage 1: LLM slide spec ────────────────────────────────────────

_SLIDES_SYSTEM = """\
You are a Chief Supply Chain Officer designing a Global Fulfilment Dashboard
presentation for a board-level audience.  You are given JSON extracted from
the Dashboard_Update worksheet.  Generate a COMPLETE slide specification as JSON.

═══ DASHBOARD TABLE DESIGN ═══

The dashboard is a SINGLE TABLE containing ALL product groups.  Each row represents
one plant-level risk item.  Rows belonging to the same product group are listed
consecutively (they will be visually grouped via vertical merge on the slide).

COLUMNS (in order):
  product_group  — "Name\\n(code)" format, e.g. "Compressors\\n(11)".
                   Every row in the same group must have the IDENTICAL value.
  plant          — Plant location code, ≤ 15 chars
  customer       — Customer or channel name, ≤ 40 chars.  May be empty.
  kb_coverage    — KB coverage to customer as compact text (e.g. "CW17/18\\nCW21").
                   Combine coverage-without and coverage-with values on separate lines.
                   If only one value, use a single line.  Use "CWnn" format.
  cw_colors      — Object mapping each of the 6 CW column numbers (as strings)
                   to a RAG colour: "GREEN", "AMBER", "RED", or "GREY".
  quarter_color  — Single RAG colour for the quarter column.
  supplier       — Supplier name or description, ≤ 60 chars
  comment        — Action / comment text, ≤ 90 chars
  fm_detail      — Force Majeure detail letter status, ≤ 30 chars

CW COLUMNS:
  Exactly 6 consecutive calendar-week integers starting from the current CW.
  Handle year-end wrap (e.g. CW50,51,52,1,2,3).

QUARTER:
  The quarter that follows the 6-week CW window (e.g. "Q2" or "Q3/2026").

RAG COLOUR COMPUTATION (per row, per CW column):
  Given coverage_without_mitigation_cw (W) and coverage_with_mitigation_cw (M):
    • "GREEN" if CW ≤ W   (supply secured without mitigation)
    • "AMBER" if W < CW ≤ M   (supply depends on mitigation)
    • "RED"   if CW > M   (no supply plan in place)
    • If M is null but W is set: "GREEN" if CW ≤ W, else "RED"
    • If both are null: "GREY"
  quarter_color = worst RAG across ALL weeks of that quarter
  (RED > AMBER > GREEN > GREY).

═══ OUTPUT FORMAT ═══

Respond with ONLY valid JSON — no markdown fences, no explanation text.

{{
  "presentation_title": "Global Fulfilment Dashboard",
  "current_cw": "CW13/2026",
  "cw_columns": [13, 14, 15, 16, 17, 18],
  "quarter_label": "Q2",
  "overall_risk": "HIGH",
  "stats": {{
    "total_items": 24,
    "red_count": 12,
    "amber_count": 5,
    "green_count": 7
  }},
  "rows": [
    {{
      "product_group": "Compressors\\n(11)",
      "plant": "LIS",
      "customer": "",
      "kb_coverage": "CW17/18\\nCW21\\nCW17",
      "cw_colors": {{"13":"GREEN","14":"GREEN","15":"GREEN","16":"GREEN","17":"GREEN","18":"RED"}},
      "quarter_color": "RED",
      "supplier": "Supplier Corp.",
      "comment": "Dual source activation by CW14",
      "fm_detail": "sent"
    }},
    {{
      "product_group": "Compressors\\n(11)",
      "plant": "LIB",
      "customer": "BMW Group",
      "kb_coverage": "CW21",
      "cw_colors": {{"13":"GREEN","14":"GREEN","15":"GREEN","16":"GREEN","17":"GREEN","18":"GREEN"}},
      "quarter_color": "RED",
      "supplier": "",
      "comment": "",
      "fm_detail": ""
    }}
  ]
}}

IMPORTANT:
  • EVERY data row from the extraction must appear — skip nothing.
  • Rows of the same product group MUST be listed consecutively with the
    EXACT SAME product_group string value (the renderer uses this to merge cells).
  • overall_risk = worst single RAG across ALL rows and CW columns.
    Map: any RED → "CRITICAL" or "HIGH", all AMBER → "MEDIUM", all GREEN → "LOW".
  • stats.red_count / amber_count / green_count = number of rows whose worst CW
    colour is RED / AMBER / GREEN respectively.
{glossary_block}"""

_SLIDES_USER = """\
Today: {today}
Current calendar week: {current_cw}

Generate the complete dashboard slide specification from this extracted data.
Use exactly 6 CW columns starting from the current week number.

{extracted_json}"""


async def llm_generate_slide_spec(
    extracted: dict,
    llm_config: dict,
    session_id: str,
    glossary_context: str = "",
) -> dict:
    """
    Stage 1: LLM generates a complete slide specification from extracted data.

    Returns a dict with a "rows" list.  Falls back to a deterministic spec on error.
    """
    llm = _create_llm(llm_config, max_tokens=8192)
    t0 = time.time()
    today = date.today().isoformat()
    current_cw = extracted.get("current_cw", "CW??/????")

    glossary_block = (
        f"\n\nCOMPANY GLOSSARY:\n{glossary_context}" if glossary_context else ""
    )

    clean = {k: v for k, v in extracted.items() if not k.startswith("_")}
    extracted_json = json.dumps(clean, indent=2, ensure_ascii=False)

    messages = [
        SystemMessage(content=_SLIDES_SYSTEM.format(
            today=today,
            glossary_block=glossary_block,
        )),
        HumanMessage(content=_SLIDES_USER.format(
            today=today,
            current_cw=current_cw,
            extracted_json=extracted_json,
        )),
    ]

    try:
        response = await llm.ainvoke(messages)
        raw = response.content.strip()

        if raw.startswith("```"):
            raw = "\n".join(raw.split("\n")[1:])
        if raw.endswith("```"):
            raw = "\n".join(raw.split("\n")[:-1])

        spec: dict = json.loads(raw.strip())

        usage = response.response_metadata.get("token_usage", {})
        log_tokens(session_id, "gfd_llm_slide_spec", usage, llm_config.get("azure_deployment", ""))

        n_rows = len(spec.get("rows", []))
        duration = (time.time() - t0) * 1000
        log_trace(
            session_id, "gfd_llm_slide_spec",
            f"Input: {len(extracted.get('product_groups', []))} product groups",
            f"Generated spec with {n_rows} rows",
            duration,
        )
        return spec

    except json.JSONDecodeError as exc:
        duration = (time.time() - t0) * 1000
        log_trace(session_id, "gfd_llm_slide_spec",
                  "Generating slide spec",
                  f"JSON PARSE ERROR — using deterministic fallback: {str(exc)[:100]}",
                  duration, {"error": True, "fallback": True})
        return _deterministic_fallback_spec(extracted)

    except Exception as exc:
        duration = (time.time() - t0) * 1000
        log_trace(session_id, "gfd_llm_slide_spec",
                  "Generating slide spec",
                  f"ERROR — using deterministic fallback: {str(exc)[:100]}",
                  duration, {"error": True, "fallback": True})
        return _deterministic_fallback_spec(extracted)


# ─── Deterministic fallback spec ────────────────────────────────────

def _cw_color(cw_num: int, cov_wo: int | None, cov_w: int | None) -> str:
    if cov_wo is None and cov_w is None:
        return "GREY"
    boundary_w = cov_w or cov_wo
    boundary_wo = cov_wo or cov_w
    if cw_num <= (boundary_wo or 0):
        return "GREEN"
    if cw_num <= (boundary_w or 0):
        return "AMBER"
    return "RED"


def _worst_rag(colors: list[str]) -> str:
    rank = {"GREY": 0, "GREEN": 1, "AMBER": 2, "RED": 3}
    return max(colors, key=lambda c: rank.get(c, 0), default="GREY")


def _deterministic_fallback_spec(extracted: dict) -> dict:
    """
    Fallback: compute slide spec without LLM using arithmetic RAG logic.
    Produces a correct dashboard but with less polished text than the LLM version.
    """
    import re as _re

    current_cw_str = extracted.get("current_cw", "CW1/2026")
    m = _re.match(r"CW(\d+)/(\d+)", current_cw_str)
    if m:
        cw_start, cw_year = int(m.group(1)), int(m.group(2))
    else:
        iso = datetime.now().isocalendar()
        cw_start, cw_year = iso.week, iso.year

    # 6 consecutive CW numbers with year-end wrap
    cw_columns: list[int] = []
    for i in range(6):
        w = ((cw_start - 1 + i) % 52) + 1
        cw_columns.append(w)

    # Quarter label
    last_cw = cw_columns[-1]
    q_num = (last_cw - 1) // 13 + 2
    q_label = f"Q{min(q_num, 4)}"

    all_rows: list[dict] = []
    red_count = 0
    amber_count = 0
    green_count = 0

    for pg in extracted.get("product_groups", []):
        code = pg.get("product_family_code", "")
        desc = pg.get("product_family_desc", "Unknown")
        pg_label = f"{desc}\n({code})" if code else desc

        for row in pg.get("rows", []):
            cov_wo: int | None = row.get("coverage_without_mitigation_cw")
            cov_w:  int | None = row.get("coverage_with_mitigation_cw")

            cw_colors = {str(cw): _cw_color(cw, cov_wo, cov_w) for cw in cw_columns}
            all_colors = list(cw_colors.values())

            # Quarter color
            quarter_weeks = range((q_num - 2) * 13 + 1, (q_num - 1) * 13 + 1)
            quarter_colors = [_cw_color(w, cov_wo, cov_w) for w in quarter_weeks if 1 <= w <= 52]
            quarter_color = _worst_rag(quarter_colors) if quarter_colors else _worst_rag(all_colors)

            # KB coverage text
            parts = []
            if cov_wo is not None:
                parts.append(f"CW{cov_wo}")
            if cov_w is not None and cov_w != cov_wo:
                parts.append(f"CW{cov_w}")
            kb_coverage = "\n".join(parts) if parts else ""

            worst = _worst_rag(all_colors)
            if worst == "RED":
                red_count += 1
            elif worst == "AMBER":
                amber_count += 1
            else:
                green_count += 1

            all_rows.append({
                "product_group":  pg_label,
                "plant":          (row.get("plant_location") or "")[:15],
                "customer":       (row.get("customer_affected") or "")[:40],
                "kb_coverage":    kb_coverage,
                "cw_colors":      cw_colors,
                "quarter_color":  quarter_color,
                "supplier":       (row.get("supplier_text") or "")[:60],
                "comment":        (row.get("action_comment") or "")[:90],
                "fm_detail":      str(row.get("customer_informed") or "")[:30],
            })

    total = len(all_rows)
    if red_count > 0:
        overall_risk = "HIGH"
    elif amber_count > 0:
        overall_risk = "MEDIUM"
    else:
        overall_risk = "LOW"

    return {
        "presentation_title": "Global Fulfilment Dashboard",
        "current_cw":         current_cw_str,
        "cw_columns":         cw_columns,
        "quarter_label":      q_label,
        "overall_risk":       overall_risk,
        "stats": {
            "total_items":  total,
            "red_count":    red_count,
            "amber_count":  amber_count,
            "green_count":  green_count,
        },
        "_fallback":          True,
        "rows":               all_rows,
    }


# ─── Stage 2: PPTX renderer ────────────────────────────────────────

def _compute_product_group_ranges(rows: list[dict]) -> list[tuple[int, int, str]]:
    """
    Identify consecutive row spans that share the same product_group value.

    Returns a list of (start_idx, end_idx, group_label) tuples
    where start_idx and end_idx are inclusive 0-based data-row indices.
    """
    if not rows:
        return []

    ranges: list[tuple[int, int, str]] = []
    current_group = rows[0].get("product_group", "")
    start = 0

    for i in range(1, len(rows)):
        grp = rows[i].get("product_group", "")
        if grp != current_group:
            ranges.append((start, i - 1, current_group))
            current_group = grp
            start = i
    ranges.append((start, len(rows) - 1, current_group))

    return ranges


def _paginate_rows(rows: list[dict], max_per_slide: int = _MAX_ROWS_PER_SLIDE,
                   ) -> list[list[dict]]:
    """
    Split rows into pages, keeping product-group boundaries intact where
    possible.  If a single product group exceeds max_per_slide, it is
    allowed to overflow onto its own page(s).
    """
    if not rows:
        return [[]]

    groups = _compute_product_group_ranges(rows)
    pages: list[list[dict]] = []
    current_page: list[dict] = []

    for start, end, _label in groups:
        group_rows = rows[start: end + 1]
        group_size = len(group_rows)

        if len(current_page) + group_size <= max_per_slide:
            current_page.extend(group_rows)
        else:
            # Current page is non-empty and group won't fit — flush
            if current_page:
                pages.append(current_page)
                current_page = []

            # If the group itself is larger than a page, chunk it
            if group_size > max_per_slide:
                for chunk_start in range(0, group_size, max_per_slide):
                    chunk = group_rows[chunk_start: chunk_start + max_per_slide]
                    pages.append(chunk)
            else:
                current_page = list(group_rows)

    if current_page:
        pages.append(current_page)

    return pages if pages else [[]]


def _render_data_slide(prs: Presentation, layout, slide_title: str,
                       rows: list[dict], cw_columns: list[int],
                       quarter_label: str, page_num: int = 0,
                       total_pages: int = 1) -> None:
    """Render one dashboard table slide matching the corporate template."""
    slide = _add_slide(prs, layout)

    # ── Title ────────────────────────────────────────────────────────
    # Try to find an existing title placeholder from the layout first
    title_set = False
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:   # title placeholder
            shape.text = slide_title
            title_set = True
            break

    if not title_set:
        # Fallback: add a textbox in the same position as the template title
        tb = slide.shapes.add_textbox(
            int(_TITLE_LEFT), int(_TITLE_TOP),
            int(_TITLE_WIDTH), int(_TITLE_HEIGHT),
        )
        tb.fill.background()
        p = tb.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = slide_title
        run.font.name = FONT
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = _C["HDR_BLUE"]

    if not rows:
        return

    # ── Table geometry ───────────────────────────────────────────────
    n_data_rows = len(rows)
    n_total_rows = n_data_rows + 2   # 2 header rows

    # Compute available height for table
    tbl_available_h = _H - _TBL_TOP - Inches(0.15)   # small bottom margin
    data_row_h = int((tbl_available_h - _HDR_ROW_H - _HDR_SUB_H) / max(n_data_rows, 1))
    # Clamp row height to reasonable range
    data_row_h = max(data_row_h, Emu(80000))    # min ~0.088 in
    data_row_h = min(data_row_h, Emu(280000))   # max ~0.306 in

    tbl_h = _HDR_ROW_H + _HDR_SUB_H + data_row_h * n_data_rows
    tbl_w = sum(_COL_WIDTHS_EMU)

    tbl_shape = slide.shapes.add_table(
        n_total_rows, _N_COLS,
        int(_TBL_LEFT), int(_TBL_TOP),
        int(tbl_w), int(tbl_h),
    )
    tbl = tbl_shape.table

    # Disable banding (we handle our own colours)
    tblPr = tbl._tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("lastRow", "0")

    # Apply column widths
    for ci, w in enumerate(_COL_WIDTHS_EMU):
        tbl.columns[ci].width = w

    # Apply row heights
    tbl.rows[0].height = int(_HDR_ROW_H)
    tbl.rows[1].height = int(_HDR_SUB_H)
    for ri in range(2, n_total_rows):
        tbl.rows[ri].height = data_row_h

    # ── Define header labels ────────────────────────────────────────
    hdr_labels = [
        "Product\nGroup\n(PG)",
        "Plant",
        "Customer/ \nChannel (PG)",
        "KB Coverage\n to Customer\n(CW/YY)",
    ]
    hdr_right_labels = ["Supplier", "Comment", "FM Detail Letter"]

    # ── Perform ALL merges BEFORE writing any content ──────────────
    #    python-pptx merge() concatenates text from all cells in the
    #    range, so merging must happen while cells are still empty.

    # Header text columns: merge row 0 + row 1 (vertical)
    for ci in list(range(_N_FIXED_LEFT)) + list(range(_Q_COL + 1, _N_COLS)):
        _merge_cells(tbl, 0, ci, 1, ci)

    # CW super-header: merge row 0 cols 4–10 (horizontal)
    _merge_cells(tbl, 0, _CW_COL_START, 0, _Q_COL)

    # Product-group column: merge data rows that share the same group
    pg_ranges = _compute_product_group_ranges(rows)
    for start, end, _label in pg_ranges:
        if end > start:
            _merge_cells(tbl, start + 2, 0, end + 2, 0)   # +2 for header rows

    # ── Now write header content ─────────────────────────────────────
    # Left text columns (merged across rows 0–1)
    for ci, label in enumerate(hdr_labels):
        cell = tbl.cell(0, ci)
        _cell_fill(cell, _C["HDR_BLUE"])
        _cell_write_multiline(cell, label, size=7, bold=False, color=_C["WHITE"])

    # CW super-header band (merged across cols 4–10 in row 0) — green bg, no text
    cell_cw_hdr = tbl.cell(0, _CW_COL_START)
    _cell_fill(cell_cw_hdr, _C["HDR_GREEN"])

    # Right text columns (merged across rows 0–1)
    for i, label in enumerate(hdr_right_labels):
        ci = _Q_COL + 1 + i
        cell = tbl.cell(0, ci)
        _cell_fill(cell, _C["HDR_BLUE"])
        _cell_write_multiline(cell, label, size=7, bold=False, color=_C["WHITE"])

    # CW number cells (row 1, cols 4–9) — green bg, white bold number
    for i, cw_num in enumerate(cw_columns):
        ci = _CW_COL_START + i
        cell = tbl.cell(1, ci)
        _cell_fill(cell, _C["HDR_GREEN"])
        _cell_write(cell, str(cw_num), size=7, bold=True,
                    color=_C["WHITE"], align=PP_ALIGN.CENTER)

    # Quarter label cell (row 1, col 10)
    q_cell = tbl.cell(1, _Q_COL)
    _cell_fill(q_cell, _C["HDR_GREEN"])
    _cell_write(q_cell, quarter_label, size=7, bold=True,
                color=_C["WHITE"], align=PP_ALIGN.CENTER)

    # ── Data rows ────────────────────────────────────────────────────
    for ri_data, row in enumerate(rows):
        ri = ri_data + 2   # offset for 2 header rows

        # Plant (may contain newlines e.g. "ALD\nBX")
        _cell_write_multiline(tbl.cell(ri, 1), row.get("plant", ""), size=6)

        # Customer / Channel
        _cell_write(tbl.cell(ri, 2), row.get("customer", ""), size=6)

        # KB Coverage (multi-line)
        _cell_write_multiline(tbl.cell(ri, 3), row.get("kb_coverage", ""), size=6)

        # CW coloured cells
        cw_colors: dict = row.get("cw_colors", {})
        for i, cw_num in enumerate(cw_columns):
            ci = _CW_COL_START + i
            rag = str(cw_colors.get(str(cw_num), "GREY")).upper()
            cell = tbl.cell(ri, ci)
            _cell_fill(cell, _RAG_BG.get(rag, _C["GREY"]))
            _cell_write(cell, "", size=5)   # empty — colour carries the meaning

        # Quarter column
        q_rag = str(row.get("quarter_color", "GREY")).upper()
        q_cell = tbl.cell(ri, _Q_COL)
        _cell_fill(q_cell, _RAG_BG.get(q_rag, _C["GREY"]))
        _cell_write(q_cell, "", size=5)

        # Supplier
        _cell_write(tbl.cell(ri, 11), row.get("supplier", ""), size=6)

        # Comment
        _cell_write(tbl.cell(ri, 12), row.get("comment", ""), size=6)

        # FM Detail Letter
        _cell_write(tbl.cell(ri, 13), row.get("fm_detail", ""), size=6)

    # ── Product-group merged cell: write label + amber fill ──────────
    #    After merge, tbl.cell(first_row, 0) refers to the merged cell.
    for start, _end, label in pg_ranges:
        cell_pg = tbl.cell(start + 2, 0)       # +2 for header rows
        _cell_fill(cell_pg, _C["AMBER"])
        _cell_write_multiline(cell_pg, label, size=6, bold=True)


def render_pptx_from_spec(slide_spec: dict,
                          output_path: str | None = None,
                          template_path: str | Path | None = None,
                          ) -> io.BytesIO:
    """
    Stage 2 (deterministic renderer): Convert an LLM-generated slide spec dict
    into a python-pptx Presentation matching the corporate GFD template.

    Parameters
    ----------
    slide_spec    : dict produced by llm_generate_slide_spec (or the fallback)
    output_path   : optional path to save the file to disk
    template_path : optional path to the corporate template PPTX

    Returns
    -------
    io.BytesIO buffer (seeked to 0) containing the .pptx bytes
    """
    prs, layout = _load_base_presentation(template_path)

    title         = slide_spec.get("presentation_title", "Global Fulfilment Dashboard")
    cw_columns    = slide_spec.get("cw_columns", [])
    quarter_label = slide_spec.get("quarter_label", "Q?")
    all_rows      = slide_spec.get("rows", [])

    # Paginate rows across slides
    pages = _paginate_rows(all_rows)
    total_pages = len(pages)

    for page_idx, page_rows in enumerate(pages):
        slide_title = title
        if total_pages > 1:
            slide_title = f"{title}  ({page_idx + 1}/{total_pages})"

        _render_data_slide(
            prs, layout,
            slide_title=slide_title,
            rows=page_rows,
            cw_columns=cw_columns,
            quarter_label=quarter_label,
            page_num=page_idx,
            total_pages=total_pages,
        )

    # Update spec with computed slide count for callers
    slide_spec["slide_count"] = total_pages

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    if output_path:
        with open(output_path, "wb") as fh:
            fh.write(buf.read())
        buf.seek(0)

    return buf


# ─── Combined entry point ───────────────────────────────────────────

async def generate_gfd_dashboard(
    extracted_data: dict,
    llm_config: dict,
    session_id: str,
    output_path: str | None = None,
    glossary_context: str = "",
) -> tuple[io.BytesIO, dict]:
    """
    Full GFD dashboard generation pipeline.

    1. LLM generates a complete slide specification from extracted_data.
    2. Thin PPTX renderer converts the spec to a python-pptx file.

    Parameters
    ----------
    extracted_data : output of parse_gfd_with_llm()
    llm_config     : Azure OpenAI config dict
    session_id     : for token / trace logging
    output_path    : optional path to save the .pptx file
    glossary_context : rendered glossary string for LLM context

    Returns
    -------
    (buf: BytesIO, slide_spec: dict)
    """
    spec = await llm_generate_slide_spec(
        extracted_data, llm_config, session_id, glossary_context
    )
    buf = render_pptx_from_spec(spec, output_path=output_path)
    return buf, spec
