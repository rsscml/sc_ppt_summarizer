"""
GFD Slide Generator Module
============================
Generates PowerPoint slides for the Global Fulfilment Dashboard.

Takes parsed data from `gfd_excel_parser.parse_dashboard_update()` and
produces a dense, color-coded dashboard table across one or more slides,
matching the operational reporting format.

The time-status grid uses **Calendar Weeks (CW)** starting from the
current week and extending 12 weeks forward. RAG status per cell is
derived from the two coverage boundary fields:

  RED   = CW is before coverage_without_mitigation (uncovered)
  AMBER = CW is between coverage w/o and w/ mitigation (at risk)
  GREEN = CW is at or past coverage_with_mitigation (covered)

Layout per slide:
  ┌──────────────────────────────────────────────────────────────────────────┐
  │  Global Fulfilment Dashboard                              [CW13/2026]  │
  ├───────┬─────┬─────────┬──────┬───┬───┬───┬───┬···┬───┬────┬──────┬────┬───┤
  │ PG    │Plant│Customer │Cover.│13 │14 │15 │16 │   │24 │ Q2 │Suppl.│Act.│FM │
  │(merge)│     │         │      │RAG│RAG│RAG│RAG│   │RAG│RAG │      │    │   │
  └───────┴─────┴─────────┴──────┴───┴───┴───┴───┴···┴───┴────┴──────┴────┴───┘
"""

import io
from pathlib import Path
from datetime import datetime, timezone
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from gfd_excel_parser import (
    get_current_cw, parse_cw, cw_to_absolute, absolute_to_cw,
)


# ─── Constants ───────────────────────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

TABLE_LEFT = Inches(0.25)
TABLE_TOP = Inches(0.75)
TABLE_WIDTH = Inches(12.83)

MAX_ROWS_PER_SLIDE = 18
CW_GRID_SPAN = 12        # number of calendar-week columns to show

# Colors
CLR_TITLE_BG = RGBColor(0x1E, 0x27, 0x61)
CLR_TITLE_FG = RGBColor(0xFF, 0xFF, 0xFF)
CLR_HEADER_BG = RGBColor(0x2C, 0x52, 0x82)
CLR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
CLR_PG_BG = RGBColor(0xFF, 0xD9, 0x66)
CLR_PG_FG = RGBColor(0x1A, 0x1A, 0x1A)
CLR_ROW_EVEN = RGBColor(0xF8, 0xF8, 0xFA)
CLR_ROW_ODD = RGBColor(0xFF, 0xFF, 0xFF)
CLR_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
CLR_MUTED_TEXT = RGBColor(0x66, 0x66, 0x66)
CLR_CW_CURRENT_HEADER = RGBColor(0x1E, 0x27, 0x61)  # highlight current CW header

RAG_COLORS = {
    "GREEN": RGBColor(0x00, 0xB0, 0x50),
    "AMBER": RGBColor(0xFF, 0xC0, 0x00),
    "RED":   RGBColor(0xFF, 0x00, 0x00),
}
RAG_TEXT_COLORS = {
    "GREEN": RGBColor(0xFF, 0xFF, 0xFF),
    "AMBER": RGBColor(0x1A, 0x1A, 0x1A),
    "RED":   RGBColor(0xFF, 0xFF, 0xFF),
}

FONT_NAME = "Calibri"
FONT_SIZE_TITLE = Pt(16)
FONT_SIZE_HEADER = Pt(7)
FONT_SIZE_CELL = Pt(6.5)
FONT_SIZE_PG = Pt(7)
FONT_SIZE_SUBTITLE = Pt(8)
FONT_SIZE_CW_HEADER = Pt(6.5)

# Fixed left-side columns
FIXED_COLUMNS = [
    ("product_group",   "Product Group (PG)", 1.45, PP_ALIGN.LEFT),
    ("plant",           "Plant",              0.55, PP_ALIGN.CENTER),
    ("customer",        "Customer / Channel", 1.20, PP_ALIGN.LEFT),
    ("coverage",        "KB Coverage\n(to Customer)", 0.80, PP_ALIGN.CENTER),
]

# Fixed right-side columns
TRAILING_COLUMNS = [
    ("supplier",  "Supplier",          1.40, PP_ALIGN.LEFT),
    ("comment",   "Action / Comment",  3.30, PP_ALIGN.LEFT),
    ("fm_detail", "FM Detail",         0.70, PP_ALIGN.CENTER),
]

CW_COL_WIDTH = 0.37  # inches per CW column — narrower to fit 12 columns
Q_COL_WIDTH = 0.50   # slightly wider for the quarter summary column
CLR_Q_HEADER = RGBColor(0x1A, 0x3C, 0x6E)  # darker blue for quarter header


# ─── Quarter Helpers ─────────────────────────────────────────────────

def _cw_to_quarter(year: int, week: int) -> tuple[int, int]:
    """Map (year, week) to (year, quarter). Q1=CW1-13, Q2=CW14-26, Q3=CW27-39, Q4=CW40-52."""
    q = (week - 1) // 13 + 1
    return (year, min(q, 4))


def _next_quarter(year: int, quarter: int) -> tuple[int, int]:
    """Return the (year, quarter) immediately following the given one."""
    if quarter >= 4:
        return (year + 1, 1)
    return (year, quarter + 1)


def _quarter_cw_range(year: int, quarter: int) -> list[tuple[int, int]]:
    """Return all (year, week) tuples belonging to a given quarter."""
    start_cw = (quarter - 1) * 13 + 1
    end_cw = quarter * 13
    return [(year, w) for w in range(start_cw, min(end_cw, 52) + 1)]


def _compute_quarter_rag(quarter_weeks: list[tuple[int, int]],
                         cov_wo: Optional[tuple[int, int]],
                         cov_w: Optional[tuple[int, int]]) -> str:
    """
    Compute worst-case RAG across all weeks in a quarter.
    Priority: RED > AMBER > GREEN > NONE.
    """
    worst = "NONE"
    rank = {"NONE": 0, "GREEN": 1, "AMBER": 2, "RED": 3}
    for cw in quarter_weeks:
        rag = _compute_cw_rag(cw, cov_wo, cov_w)
        if rank.get(rag, 0) > rank.get(worst, 0):
            worst = rag
    return worst


# ─── CW Grid Generation ─────────────────────────────────────────────

def _generate_cw_grid(start_cw: tuple[int, int], span: int = CW_GRID_SPAN
                      ) -> list[tuple[int, int]]:
    """
    Generate a list of (year, week) tuples starting from start_cw,
    spanning `span` weeks forward.
    """
    grid = []
    abs_start = cw_to_absolute(start_cw[0], start_cw[1])
    for offset in range(span):
        abs_week = abs_start + offset
        y, w = absolute_to_cw(abs_week)
        # Handle week overflow: ISO years have 52 or 53 weeks
        # Our simple model uses 53 slots per year; in practice
        # weeks 1–52 cover all cases; week 53 is rare.
        if w > 52:
            y += 1
            w = w - 52
        grid.append((y, w))
    return grid


def _compute_cw_rag(cw: tuple[int, int],
                    cov_wo: Optional[tuple[int, int]],
                    cov_w: Optional[tuple[int, int]]) -> str:
    """
    Compute RAG status for a single CW cell.

    "Coverage w/o mitigation = CW18" means supply is secured through CW18
    without any mitigation actions. "Coverage w/ mitigation = CW22" means
    supply is secured through CW22 if mitigations succeed.

    Logic:
      GREEN — CW is within coverage_without_mitigation horizon (secured)
      AMBER — CW is beyond w/o mitigation but within w/ mitigation (at risk)
      RED   — CW is beyond all coverage (uncovered)

    Edge cases:
      - If only cov_wo: GREEN up to it, RED after
      - If only cov_w:  GREEN up to it, RED after
      - If neither:     no RAG (return "NONE")
    """
    cw_abs = cw_to_absolute(cw[0], cw[1])

    if cov_wo and cov_w:
        wo_abs = cw_to_absolute(cov_wo[0], cov_wo[1])
        w_abs  = cw_to_absolute(cov_w[0], cov_w[1])
        if cw_abs <= wo_abs:
            return "GREEN"
        elif cw_abs <= w_abs:
            return "AMBER"
        else:
            return "RED"

    boundary = cov_w or cov_wo
    if boundary:
        b_abs = cw_to_absolute(boundary[0], boundary[1])
        return "GREEN" if cw_abs <= b_abs else "RED"

    return "NONE"


# ─── Helpers ─────────────────────────────────────────────────────────

def _truncate(text: Optional[str], max_len: int = 120) -> str:
    if not text:
        return ""
    s = str(text).strip()
    return s if len(s) <= max_len else s[:max_len - 1] + "…"


def _set_cell_text(cell, text: str, font_size=FONT_SIZE_CELL,
                   bold: bool = False, color: RGBColor = CLR_DARK_TEXT,
                   alignment: int = PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = Pt(8)

    for run in p.runs:
        run.font.name = FONT_NAME
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = color


def _set_cell_fill(cell, color: RGBColor):
    cell_fill = cell.fill
    cell_fill.solid()
    cell_fill.fore_color.rgb = color


# ─── Slide Row Preparation ───────────────────────────────────────────

def _prepare_slide_rows(product_groups: list[dict]) -> list[dict]:
    """
    Flatten product groups into slide rows with PG merge metadata.

    Prefers LLM-interpreted fields (llm_*) when present, falls back to raw.
    CW coverage boundaries are always taken from the parser (deterministic).
    """
    slide_rows = []

    for pg in product_groups:
        desc = pg.get("product_family_desc") or "Unknown"
        code = pg.get("product_family_code") or ""
        # Use LLM headline if available, otherwise construct from desc+code
        pg_label = pg.get("pg_headline") or (f"{desc} ({code})" if code else desc)
        pg_row_count = len(pg["rows"])

        for i, row in enumerate(pg["rows"]):
            slide_rows.append({
                "pg_label": pg_label,
                "pg_span": pg_row_count if i == 0 else 0,
                "plant": _truncate(row.get("plant_location"), 15),
                "customer": _truncate(
                    row.get("llm_customer_summary") or
                    row.get("customer_affected") or
                    _derive_affected_customers(row), 40
                ),
                "coverage": _truncate(
                    row.get("coverage_with_mitigation") or
                    row.get("coverage_without_mitigation"), 20
                ),
                "coverage_cw_wo": row.get("coverage_cw_without_mitigation"),
                "coverage_cw_w":  row.get("coverage_cw_with_mitigation"),
                "supplier": _truncate(
                    row.get("llm_supplier") or
                    row.get("supplier_text"), 35
                ),
                "comment": _truncate(
                    row.get("llm_action") or
                    row.get("action_comment"), 120
                ),
                "fm_detail": _truncate(
                    row.get("llm_fm_status") or _format_fm(row), 15
                ),
                "risk_level": row.get("llm_risk_level", ""),
            })

    return slide_rows


def _derive_affected_customers(row: dict) -> str:
    impacts = row.get("customer_impact", {})
    affected = [k for k, v in impacts.items() if v]
    if not affected:
        return ""
    if len(affected) <= 3:
        return ", ".join(affected)
    return ", ".join(affected[:3]) + f" +{len(affected) - 3}"


def _format_fm(row: dict) -> str:
    informed = row.get("customer_informed")
    if informed is True:
        return "Yes"
    if informed is False:
        return "No"
    return _truncate(str(row.get("customer_informed", "")), 15)


def _paginate(slide_rows: list[dict], max_per_slide: int = MAX_ROWS_PER_SLIDE
              ) -> list[list[dict]]:
    """Split rows into pages without breaking product-group spans."""
    pages: list[list[dict]] = []
    current_page: list[dict] = []
    i = 0

    while i < len(slide_rows):
        row = slide_rows[i]
        span = row["pg_span"]

        if span > 0:
            group_size = span
            group_rows = slide_rows[i:i + group_size]

            if len(current_page) + group_size > max_per_slide:
                if current_page:
                    pages.append(current_page)
                    current_page = []

                if group_size > max_per_slide:
                    while group_rows:
                        chunk = group_rows[:max_per_slide]
                        chunk[0] = {**chunk[0], "pg_span": len(chunk)}
                        pages.append(chunk)
                        group_rows = group_rows[max_per_slide:]
                        if group_rows:
                            group_rows[0] = {**group_rows[0], "pg_span": len(group_rows)}
                    i += group_size
                    continue

            current_page.extend(group_rows)
            i += group_size
        else:
            current_page.append(row)
            i += 1

        if len(current_page) >= max_per_slide:
            pages.append(current_page)
            current_page = []

    if current_page:
        pages.append(current_page)

    return pages if pages else [[]]


# ─── Column Spec Builder ────────────────────────────────────────────

def _build_column_spec(cw_grid: list[tuple[int, int]],
                       next_q: tuple[int, int]) -> list[tuple]:
    """
    Build the full column specification for the slide table.
    CW columns + 1 next-quarter column are inserted between FIXED and TRAILING.
    """
    fixed_width = sum(c[2] for c in FIXED_COLUMNS)
    cw_width = len(cw_grid) * CW_COL_WIDTH + Q_COL_WIDTH
    trailing_width = sum(c[2] for c in TRAILING_COLUMNS)
    total = fixed_width + cw_width + trailing_width
    available = TABLE_WIDTH / Inches(1)

    if total > available:
        scale = (available - fixed_width - cw_width) / trailing_width
        scaled_trailing = [(c[0], c[1], c[2] * scale, c[3]) for c in TRAILING_COLUMNS]
    else:
        scaled_trailing = list(TRAILING_COLUMNS)

    spec = list(FIXED_COLUMNS)

    # Individual CW columns
    for (year, week) in cw_grid:
        label = str(week)
        spec.append((f"cw_{year}_{week}", label, CW_COL_WIDTH, PP_ALIGN.CENTER))

    # Next-quarter summary column
    q_label = f"Q{next_q[1]}"
    spec.append((f"q_{next_q[0]}_{next_q[1]}", q_label, Q_COL_WIDTH, PP_ALIGN.CENTER))

    spec.extend(scaled_trailing)
    return spec


# ─── Title Bar ───────────────────────────────────────────────────────

def _add_title_bar(slide, page_num: int, total_pages: int,
                   current_cw: tuple[int, int]):
    title_box = slide.shapes.add_textbox(
        TABLE_LEFT, Inches(0.15), TABLE_WIDTH, Inches(0.50)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    run = p.add_run()
    run.text = "Global Fulfilment Dashboard"
    run.font.name = FONT_NAME
    run.font.size = FONT_SIZE_TITLE
    run.font.bold = True
    run.font.color.rgb = CLR_TITLE_BG

    now = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    page_text = f"CW{current_cw[1]}/{current_cw[0]}  •  Page {page_num}/{total_pages}  •  {now}"

    info_box = slide.shapes.add_textbox(
        Inches(9.0), Inches(0.20), Inches(4.0), Inches(0.40)
    )
    tf2 = info_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = page_text
    run2.font.name = FONT_NAME
    run2.font.size = FONT_SIZE_SUBTITLE
    run2.font.color.rgb = CLR_MUTED_TEXT


# ─── Table Builder ───────────────────────────────────────────────────

def _build_table_on_slide(slide, page_rows: list[dict], col_spec: list[tuple],
                          cw_grid: list[tuple[int, int]],
                          current_cw: tuple[int, int]):
    """Build the dense dashboard table with CW-based RAG grid."""
    n_cols = len(col_spec)
    n_rows = len(page_rows) + 1  # +1 header

    available_height = SLIDE_HEIGHT - TABLE_TOP - Inches(0.25)
    natural_height = available_height / n_rows
    row_height = max(Inches(0.28), min(natural_height, Inches(0.42)))
    table_height = row_height * n_rows

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, TABLE_LEFT, TABLE_TOP, TABLE_WIDTH, table_height
    )
    table = table_shape.table

    for col_idx, (_, _, width_in, _) in enumerate(col_spec):
        table.columns[col_idx].width = Inches(width_in)

    # ── Header row ────────────────────────────────────────────────
    for col_idx, (key, label, _, align) in enumerate(col_spec):
        cell = table.cell(0, col_idx)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if key.startswith("cw_"):
            # CW header — highlight the current-week column
            parts = key.split("_")
            cw_y, cw_w = int(parts[1]), int(parts[2])
            is_current = (cw_y == current_cw[0] and cw_w == current_cw[1])

            if is_current:
                _set_cell_fill(cell, CLR_CW_CURRENT_HEADER)
            else:
                _set_cell_fill(cell, CLR_HEADER_BG)
            _set_cell_text(cell, label, FONT_SIZE_CW_HEADER, bold=True,
                           color=CLR_HEADER_FG, alignment=PP_ALIGN.CENTER)
        elif key.startswith("q_"):
            # Quarter summary header — visually distinct
            _set_cell_fill(cell, CLR_Q_HEADER)
            _set_cell_text(cell, label, FONT_SIZE_HEADER, bold=True,
                           color=CLR_HEADER_FG, alignment=PP_ALIGN.CENTER)
        else:
            _set_cell_fill(cell, CLR_HEADER_BG)
            _set_cell_text(cell, label, FONT_SIZE_HEADER, bold=True,
                           color=CLR_HEADER_FG, alignment=PP_ALIGN.CENTER)

    # ── Data rows ─────────────────────────────────────────────────
    merge_tracker: dict[str, dict] = {}

    for row_idx, row_data in enumerate(page_rows):
        table_row = row_idx + 1
        is_even = row_idx % 2 == 0
        bg_color = CLR_ROW_EVEN if is_even else CLR_ROW_ODD

        cov_wo = row_data.get("coverage_cw_wo")
        cov_w  = row_data.get("coverage_cw_w")

        for col_idx, (key, label, _, align) in enumerate(col_spec):
            cell = table.cell(table_row, col_idx)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if key == "product_group":
                pg_label = row_data["pg_label"]
                span = row_data["pg_span"]
                if span > 0:
                    _set_cell_fill(cell, CLR_PG_BG)
                    _set_cell_text(cell, pg_label, FONT_SIZE_PG, bold=True,
                                   color=CLR_PG_FG, alignment=PP_ALIGN.LEFT)
                    merge_tracker[pg_label] = {"start_row": table_row, "span": span}
                else:
                    _set_cell_fill(cell, CLR_PG_BG)
                    _set_cell_text(cell, "", FONT_SIZE_PG)

            elif key == "plant":
                _set_cell_fill(cell, bg_color)
                _set_cell_text(cell, row_data.get("plant", ""),
                               alignment=PP_ALIGN.CENTER)

            elif key == "customer":
                _set_cell_fill(cell, bg_color)
                _set_cell_text(cell, row_data.get("customer", ""),
                               alignment=PP_ALIGN.LEFT)

            elif key == "coverage":
                _set_cell_fill(cell, bg_color)
                _set_cell_text(cell, row_data.get("coverage", ""),
                               alignment=PP_ALIGN.CENTER)

            elif key.startswith("cw_"):
                # ── CW RAG cell ──────────────────────────────────
                parts = key.split("_")
                cw_y, cw_w = int(parts[1]), int(parts[2])
                rag = _compute_cw_rag((cw_y, cw_w), cov_wo, cov_w)

                if rag in RAG_COLORS:
                    _set_cell_fill(cell, RAG_COLORS[rag])
                    _set_cell_text(cell, "", FONT_SIZE_CELL,
                                   color=RAG_TEXT_COLORS.get(rag, CLR_DARK_TEXT),
                                   alignment=PP_ALIGN.CENTER)
                else:
                    _set_cell_fill(cell, bg_color)
                    _set_cell_text(cell, "", alignment=PP_ALIGN.CENTER)

            elif key.startswith("q_"):
                # ── Quarter summary RAG cell ─────────────────────
                parts = key.split("_")
                q_y, q_n = int(parts[1]), int(parts[2])
                q_weeks = _quarter_cw_range(q_y, q_n)
                rag = _compute_quarter_rag(q_weeks, cov_wo, cov_w)

                if rag in RAG_COLORS:
                    _set_cell_fill(cell, RAG_COLORS[rag])
                    _set_cell_text(cell, "", FONT_SIZE_CELL,
                                   color=RAG_TEXT_COLORS.get(rag, CLR_DARK_TEXT),
                                   alignment=PP_ALIGN.CENTER)
                else:
                    _set_cell_fill(cell, bg_color)
                    _set_cell_text(cell, "", alignment=PP_ALIGN.CENTER)

            elif key == "supplier":
                _set_cell_fill(cell, bg_color)
                _set_cell_text(cell, row_data.get("supplier", ""),
                               alignment=PP_ALIGN.LEFT)

            elif key == "comment":
                _set_cell_fill(cell, bg_color)
                _set_cell_text(cell, row_data.get("comment", ""),
                               FONT_SIZE_CELL, alignment=PP_ALIGN.LEFT)

            elif key == "fm_detail":
                _set_cell_fill(cell, bg_color)
                _set_cell_text(cell, row_data.get("fm_detail", ""),
                               alignment=PP_ALIGN.CENTER)

    # ── Merge product-group cells ────────────────────────────────
    for pg_label, info in merge_tracker.items():
        start = info["start_row"]
        span = info["span"]
        if span > 1:
            try:
                top_cell = table.cell(start, 0)
                bottom_cell = table.cell(start + span - 1, 0)
                top_cell.merge(bottom_cell)
                _set_cell_fill(top_cell, CLR_PG_BG)
                _set_cell_text(top_cell, pg_label, FONT_SIZE_PG, bold=True,
                               color=CLR_PG_FG, alignment=PP_ALIGN.LEFT)
                top_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass


# ─── Executive Overview Slide ────────────────────────────────────────

RISK_BADGE_COLORS = {
    "CRITICAL": RGBColor(0xFF, 0x00, 0x00),
    "HIGH":     RGBColor(0xFF, 0x66, 0x00),
    "MEDIUM":   RGBColor(0xFF, 0xC0, 0x00),
    "LOW":      RGBColor(0x00, 0xB0, 0x50),
}


def _add_executive_overview_slide(prs, overview: dict, current_cw: tuple[int, int]):
    """Add an executive overview slide as the first slide in the deck."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    title_text = overview.get("title", f"Fulfilment Risk Overview — CW{current_cw[1]}/{current_cw[0]}")
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = FONT_NAME
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = CLR_TITLE_BG

    overall_risk = overview.get("overall_risk", "HIGH")
    badge_color = RISK_BADGE_COLORS.get(overall_risk, RISK_BADGE_COLORS["HIGH"])
    badge = slide.shapes.add_shape(1, Inches(10.5), Inches(0.3), Inches(2.2), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = badge_color
    badge.line.fill.background()
    bp = badge.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = f"Overall: {overall_risk}"
    br.font.name = FONT_NAME
    br.font.size = Pt(12)
    br.font.bold = True
    br.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    bullets = overview.get("bullets", [])
    if bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.5), Inches(5.0))
        tf2 = bullet_box.text_frame
        tf2.word_wrap = True
        for i, bullet_text in enumerate(bullets):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.space_before = Pt(8)
            p2.space_after = Pt(4)
            marker = p2.add_run()
            marker.text = "●  "
            marker.font.name = FONT_NAME
            marker.font.size = Pt(10)
            marker.font.color.rgb = CLR_TITLE_BG
            content = p2.add_run()
            content.text = bullet_text
            content.font.name = FONT_NAME
            content.font.size = Pt(14)
            content.font.color.rgb = CLR_DARK_TEXT

    now = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    footer = slide.shapes.add_textbox(Inches(9.5), Inches(6.8), Inches(3.5), Inches(0.4))
    fp = footer.text_frame.paragraphs[0]
    fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run()
    fr.text = f"Generated {now}"
    fr.font.name = FONT_NAME
    fr.font.size = Pt(8)
    fr.font.color.rgb = CLR_MUTED_TEXT


# ─── Public API ──────────────────────────────────────────────────────

def generate_gfd_slides(product_groups: list[dict],
                        executive_overview: Optional[dict] = None,
                        output_path: Optional[str] = None,
                        max_rows_per_slide: int = MAX_ROWS_PER_SLIDE,
                        cw_span: int = CW_GRID_SPAN) -> io.BytesIO:
    """
    Generate PowerPoint slides for the Global Fulfilment Dashboard.

    Args:
        product_groups:     List of product group dicts. Accepts either:
                            - raw from parser: parsed_data["product_groups"]
                            - LLM-interpreted from gfd_agent: result["interpreted_groups"]
                            When llm_* fields are present they are preferred for text columns.
        executive_overview: Optional dict with {title, bullets, overall_risk}.
                            Produced by gfd_agent; adds an overview slide before the data.
        output_path:        Optional file path to save
        max_rows_per_slide: Max data rows per slide before pagination
        cw_span:            Number of forward-looking CW columns (default 12)

    Returns:
        io.BytesIO buffer containing the .pptx file
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    current_cw = get_current_cw()
    cw_grid = _generate_cw_grid(current_cw, cw_span)
    cur_q = _cw_to_quarter(current_cw[0], current_cw[1])
    next_q = _next_quarter(cur_q[0], cur_q[1])
    col_spec = _build_column_spec(cw_grid, next_q)

    # ── Executive overview slide (when LLM pipeline provides one) ───
    if executive_overview:
        _add_executive_overview_slide(prs, executive_overview, current_cw)

    # ── Data slides ─────────────────────────────────────────────────
    slide_rows = _prepare_slide_rows(product_groups)

    if not slide_rows:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _add_title_bar(slide, 1, 1, current_cw)
        tb = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = "No data rows found in the Dashboard_Update worksheet."
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(14)
            run.font.color.rgb = CLR_MUTED_TEXT
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        if output_path:
            with open(output_path, "wb") as f:
                f.write(buf.getvalue())
            buf.seek(0)
        return buf

    pages = _paginate(slide_rows, max_rows_per_slide)
    total_pages = len(pages)

    for page_num, page_rows in enumerate(pages, 1):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _add_title_bar(slide, page_num, total_pages, current_cw)
        _build_table_on_slide(slide, page_rows, col_spec, cw_grid, current_cw)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    if output_path:
        with open(output_path, "wb") as f:
            f.write(buf.getvalue())
        buf.seek(0)
    return buf


def generate_gfd_slides_from_file(excel_path: str,
                                  output_path: Optional[str] = None,
                                  sheet_name: str = "Dashboard_Update",
                                  max_rows_per_slide: int = MAX_ROWS_PER_SLIDE,
                                  cw_span: int = CW_GRID_SPAN
                                  ) -> dict[str, Any]:
    """
    End-to-end convenience (no LLM): parse Excel → generate slides.
    For the LLM-powered pipeline, use gfd_agent.run_gfd_pipeline() instead.
    """
    from gfd_excel_parser import parse_dashboard_update

    parsed = parse_dashboard_update(excel_path, sheet_name=sheet_name)

    buf = generate_gfd_slides(
        parsed["product_groups"],
        output_path=output_path,
        max_rows_per_slide=max_rows_per_slide,
        cw_span=cw_span,
    )

    slide_rows = _prepare_slide_rows(parsed["product_groups"])
    pages = _paginate(slide_rows, max_rows_per_slide)

    return {
        "pptx_buffer": buf,
        "output_path": output_path,
        "parsed_data": parsed,
        "slide_count": len(pages),
        "row_count": len(slide_rows),
        "warnings": parsed.get("warnings", []),
    }
